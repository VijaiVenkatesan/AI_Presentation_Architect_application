from pptx import Presentation

path = r"C:\Users\vijai.v\Downloads\presentation_20260406_131646.pptx"
prs = Presentation(path)
print(f"Total slides: {len(prs.slides)}")
for i, slide in enumerate(prs.slides):
    layout = slide.slide_layout.name
    shapes = slide.shapes
    has_chart = any(str(s.shape_type) == "MSO_SHAPE_TYPE.CHART: 3" or "3)" in str(s.shape_type) for s in shapes)
    has_table = any("19" in str(s.shape_type) for s in shapes)
    has_ph = any("PLACEHOLDER" in str(s.shape_type) for s in shapes)
    visual = " [CHART]" if has_chart else (" [TABLE]" if has_table else "")
    print(f"  Slide {i+1:02d} | {layout:<30} | shapes:{len(list(shapes)):2d}{visual}")
