from pptx import Presentation

def debug():
    try:
        prs = Presentation("Datamatics-Presentation-Template-Light-Test.pptx")
    except Exception:
        print("Cannot find Datamatics presentation.")
        return
        
    for i, slide in enumerate(prs.slide_layouts):
        print(f"Layout {i}: {slide.name}")
        for ph in slide.placeholders:
            print(f"  idx:{ph.placeholder_format.idx} type:{ph.placeholder_format.type} has_frame:{ph.has_text_frame}")

debug()
