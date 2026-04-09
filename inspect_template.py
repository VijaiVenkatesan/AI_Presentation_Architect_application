
from pptx import Presentation
import os

def inspect_layouts(template_path):
    if not os.path.exists(template_path):
        print(f"Error: Template not found at {template_path}")
        return

    prs = Presentation(template_path)
    print(f"--- Template Layout Inspection: {os.path.basename(template_path)} ---")
    
    for i, layout in enumerate(prs.slide_layouts):
        print(f"\nLayout {i}: {layout.name}")
        print(f"  Shape count: {len(layout.shapes)}")
        for j, shape in enumerate(layout.shapes):
            print(f"    Shape {j+1}: {shape.name} (Type: {shape.shape_type})")
        
        # Check if it has any placeholders
        print(f"  Placeholder count: {len(layout.placeholders)}")
        for ph in layout.placeholders:
            print(f"    PH {ph.placeholder_format.idx}: {ph.name} (Type: {ph.placeholder_format.type})")

if __name__ == "__main__":
    # Try to find the template used by the user. 
    # Based on the screenshot it says "datamatics-light". 
    # The file in assets/sample_template.pptx might be it.
    template = r"D:\Antigravity_Projects\AI_Presentation_Architecture_App\assets\Datamatics-Presentation-Template-Light-Test.pptx"
    inspect_layouts(template)
