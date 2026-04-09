
import os
import sys
import io
from pptx import Presentation

# Add path
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))

from core.ppt_generator import PptGenerator
from core.template_parser import PptxTemplateParser

def test_visuals_fix():
    template_path = r"D:\Antigravity_Projects\AI_Presentation_Architecture_App\assets\Datamatics-Presentation-Template-Light-Test.pptx"
    out_path = r"C:\Users\vijai.v\Downloads\test_visuals_fix.pptx"
    
    # 1. Mock Data matching AI schema
    slides = [
        {
            "title": "Main Title Slide",
            "subtitle": "Branding Guard & Smart Layouts",
            "bullet_points": ["Point 1: Stability", "Point 2: Fidelity"]
        },
        {
            "title": "Chart Verification",
            "subtitle": "Testing single-series bar chart",
            "bullet_points": ["This slide should have a chart below"],
            "chart_data": {
                "type": "bar",
                "title": "Sales Data",
                "categories": ["2024", "2025", "2026"],
                "values": [500, 750, 1200]
            }
        },
        {
            "title": "Table Verification",
            "subtitle": "Testing {headers, rows} format",
            "bullet_points": ["This slide should have a table below"],
            "table_data": {
                "headers": ["Quarter", "Revenue", "Growth"],
                "rows": [["Q1", "$10M", "10%"], ["Q2", "$12M", "15%"]]
            }
        },
        {
            "title": "Thank You",
            "subtitle": "End of test"
        }
    ]
    
    # 2. Parse Template
    with open(template_path, "rb") as f:
        template_bytes = f.read()
    
    parser = PptxTemplateParser(template_bytes)
    profile = parser.parse()
    
    # 3. Generate
    generator = PptGenerator(template_bytes=template_bytes, template_profile=profile)
    pptx_bytes = generator.generate(slides)
    
    with open(out_path, "wb") as f:
        f.write(pptx_bytes)
    
    print(f"Test PPTX saved to {out_path}")
    
    # 4. Verify
    try:
        prs = Presentation(out_path)
        print(f"File Readable. Slide count: {len(prs.slides)}")
        for i, slide in enumerate(prs.slides):
            print(f"Slide {i+1}: {slide.slide_layout.name}")
            shape_types = [s.shape_type for s in slide.shapes]
            print(f"  Shape types: {shape_types}")
            
            # Check for specific shapes
            if i == 1: # Chart slide
                has_chart = any("CHART" in str(st) for st in shape_types)
                print(f"  Chart Found: {has_chart}")
            if i == 2: # Table slide
                has_table = any("TABLE" in str(st) for st in shape_types)
                print(f"  Table Found: {has_table}")
                
    except Exception as e:
        print(f"Verification Error: {e}")

if __name__ == "__main__":
    test_visuals_fix()
