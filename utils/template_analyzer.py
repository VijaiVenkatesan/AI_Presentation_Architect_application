"""
Template Analyzer Module - Enhanced Version
Extracts ALL styling, layout, logos, and design elements from uploaded templates
"""

import io
import re
import copy
from typing import Dict, List, Any, Optional, Tuple
from PIL import Image
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.oxml.ns import qn
from collections import Counter
import colorsys


class TemplateAnalyzer:
    """Analyzes PowerPoint templates and extracts ALL styling information"""
    
    def __init__(self):
        self.template_data = self._get_default_template()
        self.original_pptx = None  # Store original presentation
        self.extracted_images = {}  # Store logos and images
        self.background_fill = None
        self.master_slides = []
    
    def _get_default_template(self) -> Dict:
        """Return default template settings"""
        return {
            'colors': {
                'primary': '#6366F1',
                'secondary': '#8B5CF6',
                'accent': '#EC4899',
                'background': '#0F172A',
                'text_primary': '#F8FAFC',
                'text_secondary': '#94A3B8'
            },
            'fonts': {
                'title': {'name': 'Arial', 'size': 44, 'bold': True, 'color': '#F8FAFC'},
                'subtitle': {'name': 'Arial', 'size': 28, 'bold': False, 'color': '#94A3B8'},
                'body': {'name': 'Arial', 'size': 18, 'bold': False, 'color': '#F8FAFC'},
                'caption': {'name': 'Arial', 'size': 12, 'bold': False, 'color': '#94A3B8'}
            },
            'layouts': [],
            'slide_size': {'width': 13.333, 'height': 7.5},
            'has_logo': False,
            'logo_image': None,
            'logo_position': {'left': 0.3, 'top': 0.3, 'width': 1.5, 'height': 0.75},
            'background': {
                'type': 'solid',  # solid, gradient, image
                'color': '#0F172A',
                'gradient_colors': [],
                'image': None
            },
            'header_style': {},
            'footer_style': {},
            'shape_styles': [],
            'color_scheme': [],
            'master_layouts': [],
            'use_template_file': False,
            'template_bytes': None
        }
    
    def analyze_pptx(self, pptx_file: io.BytesIO) -> Dict[str, Any]:
        """Analyze a PowerPoint file and extract ALL template information"""
        try:
            # Store original bytes for later use
            pptx_file.seek(0)
            self.template_data['template_bytes'] = pptx_file.read()
            pptx_file.seek(0)
            self.template_data['use_template_file'] = True
            
            # Parse presentation
            prs = Presentation(pptx_file)
            self.original_pptx = prs
            
            # Extract slide size
            self.template_data['slide_size'] = {
                'width': prs.slide_width.inches,
                'height': prs.slide_height.inches
            }
            
            # Extract from slide masters (most important for branding)
            self._extract_master_styles(prs)
            
            # Extract from actual slides
            colors = []
            fonts = []
            layouts = []
            
            for slide_idx, slide in enumerate(prs.slides):
                slide_analysis = self._analyze_slide_complete(slide, slide_idx)
                colors.extend(slide_analysis['colors'])
                fonts.extend(slide_analysis['fonts'])
                layouts.append(slide_analysis['layout'])
                
                # Extract logo from first slide (usually title slide)
                if slide_idx == 0:
                    self._extract_logo(slide)
                    self._extract_background(slide)
            
            # Process extracted data
            self._process_colors(colors)
            self._process_fonts(fonts)
            self.template_data['layouts'] = layouts
            
            # Extract color scheme from theme
            self._extract_theme_colors(prs)
            
            return self.template_data
            
        except Exception as e:
            print(f"Error analyzing PPTX: {e}")
            import traceback
            traceback.print_exc()
            return self.template_data
    
    def _extract_master_styles(self, prs: Presentation):
        """Extract styles from slide masters"""
        try:
            for master in prs.slide_masters:
                master_info = {
                    'name': getattr(master, 'name', 'Default'),
                    'layouts': [],
                    'background': None,
                    'shapes': []
                }
                
                # Extract master background
                if hasattr(master, 'background'):
                    bg = master.background
                    if bg.fill.type is not None:
                        master_info['background'] = self._extract_fill_info(bg.fill)
                        if master_info['background']:
                            self.template_data['background'] = master_info['background']
                
                # Extract shapes from master (logos, headers, footers)
                for shape in master.shapes:
                    shape_info = self._extract_shape_complete(shape)
                    master_info['shapes'].append(shape_info)
                    
                    # Check if it's a logo (small image in corner)
                    if shape_info.get('type') == 'picture':
                        if shape_info.get('top', 0) < 1 or shape_info.get('top', 0) > 6:
                            # Likely a logo (top or bottom)
                            self._extract_image_from_shape(shape, 'logo')
                
                # Extract layouts
                for layout in master.slide_layouts:
                    layout_info = {
                        'name': layout.name,
                        'placeholders': []
                    }
                    
                    for placeholder in layout.placeholders:
                        ph_info = {
                            'type': str(placeholder.placeholder_format.type),
                            'idx': placeholder.placeholder_format.idx,
                            'left': placeholder.left.inches if placeholder.left else 0,
                            'top': placeholder.top.inches if placeholder.top else 0,
                            'width': placeholder.width.inches if placeholder.width else 0,
                            'height': placeholder.height.inches if placeholder.height else 0
                        }
                        
                        # Extract font from placeholder
                        if hasattr(placeholder, 'text_frame'):
                            for para in placeholder.text_frame.paragraphs:
                                if para.runs:
                                    run = para.runs[0]
                                    ph_info['font'] = {
                                        'name': run.font.name,
                                        'size': run.font.size.pt if run.font.size else None,
                                        'bold': run.font.bold,
                                        'color': self._get_font_color(run.font)
                                    }
                                    break
                        
                        layout_info['placeholders'].append(ph_info)
                    
                    master_info['layouts'].append(layout_info)
                
                self.template_data['master_layouts'].append(master_info)
                
        except Exception as e:
            print(f"Error extracting master styles: {e}")
    
    def _extract_logo(self, slide):
        """Extract logo from slide"""
        try:
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    # Check if it's likely a logo (small, in corner)
                    left = shape.left.inches if shape.left else 0
                    top = shape.top.inches if shape.top else 0
                    width = shape.width.inches if shape.width else 0
                    height = shape.height.inches if shape.height else 0
                    
                    # Logo criteria: small image, positioned in corners
                    is_corner = (left < 2 or left > 10) and (top < 1.5 or top > 6)
                    is_small = width < 3 and height < 2
                    
                    if is_corner and is_small:
                        image_bytes = self._extract_image_from_shape(shape, 'logo')
                        if image_bytes:
                            self.template_data['has_logo'] = True
                            self.template_data['logo_image'] = image_bytes
                            self.template_data['logo_position'] = {
                                'left': left,
                                'top': top,
                                'width': width,
                                'height': height
                            }
                            return
        except Exception as e:
            print(f"Error extracting logo: {e}")
    
    def _extract_image_from_shape(self, shape, image_type: str) -> Optional[bytes]:
        """Extract image bytes from a picture shape"""
        try:
            if hasattr(shape, 'image'):
                image_bytes = shape.image.blob
                self.extracted_images[image_type] = image_bytes
                return image_bytes
        except Exception as e:
            print(f"Error extracting image: {e}")
        return None
    
    def _extract_background(self, slide):
        """Extract background from slide"""
        try:
            if hasattr(slide, 'background'):
                bg = slide.background
                if bg.fill.type is not None:
                    bg_info = self._extract_fill_info(bg.fill)
                    if bg_info:
                        self.template_data['background'] = bg_info
        except Exception as e:
            print(f"Error extracting background: {e}")
    
    def _extract_fill_info(self, fill) -> Optional[Dict]:
        """Extract fill information (solid, gradient, pattern)"""
        try:
            fill_info = {'type': 'solid', 'color': '#0F172A'}
            
            if fill.type is not None:
                fill_type = str(fill.type)
                
                if 'SOLID' in fill_type:
                    fill_info['type'] = 'solid'
                    if fill.fore_color and fill.fore_color.rgb:
                        fill_info['color'] = self._rgb_to_hex(fill.fore_color.rgb)
                
                elif 'GRADIENT' in fill_type:
                    fill_info['type'] = 'gradient'
                    fill_info['gradient_colors'] = []
                    try:
                        for stop in fill.gradient_stops:
                            if stop.color and stop.color.rgb:
                                fill_info['gradient_colors'].append(
                                    self._rgb_to_hex(stop.color.rgb)
                                )
                    except:
                        pass
                
                elif 'PICTURE' in fill_type or 'BACKGROUND' in fill_type:
                    fill_info['type'] = 'image'
                    # Try to extract background image
                    try:
                        if hasattr(fill, '_fill'):
                            blip = fill._fill.find(qn('a:blip'))
                            if blip is not None:
                                fill_info['has_image'] = True
                    except:
                        pass
            
            return fill_info
        except Exception as e:
            return None
    
    def _analyze_slide_complete(self, slide, slide_idx: int) -> Dict[str, Any]:
        """Analyze a single slide completely"""
        analysis = {
            'colors': [],
            'fonts': [],
            'layout': {
                'slide_index': slide_idx,
                'shapes': [],
                'has_title': False,
                'has_content': False,
                'has_chart': False,
                'has_table': False,
                'has_image': False,
                'has_logo': False
            }
        }
        
        for shape in slide.shapes:
            shape_info = self._extract_shape_complete(shape)
            analysis['layout']['shapes'].append(shape_info)
            
            # Collect colors
            if shape_info.get('fill_color'):
                analysis['colors'].append(shape_info['fill_color'])
            if shape_info.get('line_color'):
                analysis['colors'].append(shape_info['line_color'])
            if shape_info.get('text_color'):
                analysis['colors'].append(shape_info['text_color'])
            
            # Collect fonts
            if shape_info.get('font'):
                analysis['fonts'].append(shape_info['font'])
            
            # Update layout flags
            shape_type = shape_info.get('type', '')
            if shape_type == 'title':
                analysis['layout']['has_title'] = True
            elif shape_type == 'text':
                analysis['layout']['has_content'] = True
            elif shape_type == 'chart':
                analysis['layout']['has_chart'] = True
            elif shape_type == 'table':
                analysis['layout']['has_table'] = True
            elif shape_type == 'picture':
                analysis['layout']['has_image'] = True
        
        return analysis
    
    def _extract_shape_complete(self, shape) -> Dict[str, Any]:
        """Extract complete information from a shape"""
        shape_info = {
            'type': 'unknown',
            'left': shape.left.inches if shape.left else 0,
            'top': shape.top.inches if shape.top else 0,
            'width': shape.width.inches if shape.width else 0,
            'height': shape.height.inches if shape.height else 0,
            'rotation': shape.rotation if hasattr(shape, 'rotation') else 0
        }
        
        # Determine shape type
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            shape_info['type'] = 'picture'
        elif shape.shape_type == MSO_SHAPE_TYPE.CHART:
            shape_info['type'] = 'chart'
            shape_info['chart_type'] = self._get_chart_type(shape)
        elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            shape_info['type'] = 'table'
            shape_info['table_style'] = self._get_table_style(shape)
        elif hasattr(shape, 'text_frame'):
            shape_info['type'] = 'text'
            text_info = self._analyze_text_frame_complete(shape.text_frame)
            shape_info.update(text_info)
            
            # Check if it's a title placeholder
            if hasattr(shape, 'placeholder_format') and shape.placeholder_format:
                ph_type = str(shape.placeholder_format.type)
                if 'TITLE' in ph_type or 'CENTER_TITLE' in ph_type:
                    shape_info['type'] = 'title'
                elif 'SUBTITLE' in ph_type:
                    shape_info['type'] = 'subtitle'
                elif 'BODY' in ph_type:
                    shape_info['type'] = 'body'
        
        # Extract fill
        if hasattr(shape, 'fill') and shape.fill:
            fill_info = self._extract_fill_info(shape.fill)
            if fill_info:
                shape_info['fill'] = fill_info
                shape_info['fill_color'] = fill_info.get('color')
        
        # Extract line/border
        if hasattr(shape, 'line') and shape.line:
            try:
                if shape.line.color and shape.line.color.rgb:
                    shape_info['line_color'] = self._rgb_to_hex(shape.line.color.rgb)
                shape_info['line_width'] = shape.line.width.pt if shape.line.width else 0
            except:
                pass
        
        return shape_info
    
    def _analyze_text_frame_complete(self, text_frame) -> Dict[str, Any]:
        """Analyze text frame completely"""
        result = {
            'paragraphs': [],
            'margin_left': text_frame.margin_left.inches if text_frame.margin_left else 0,
            'margin_right': text_frame.margin_right.inches if text_frame.margin_right else 0,
            'margin_top': text_frame.margin_top.inches if text_frame.margin_top else 0,
            'margin_bottom': text_frame.margin_bottom.inches if text_frame.margin_bottom else 0
        }
        
        for para in text_frame.paragraphs:
            para_info = {
                'alignment': str(para.alignment) if para.alignment else 'LEFT',
                'level': para.level,
                'space_before': para.space_before.pt if para.space_before else 0,
                'space_after': para.space_after.pt if para.space_after else 0,
                'runs': []
            }
            
            for run in para.runs:
                run_info = {
                    'text': run.text,
                    'font': {
                        'name': run.font.name,
                        'size': run.font.size.pt if run.font.size else 18,
                        'bold': run.font.bold if run.font.bold is not None else False,
                        'italic': run.font.italic if run.font.italic is not None else False,
                        'underline': run.font.underline if run.font.underline is not None else False,
                        'color': self._get_font_color(run.font)
                    }
                }
                para_info['runs'].append(run_info)
                
                # Store first font found for this text frame
                if 'font' not in result:
                    result['font'] = run_info['font']
                    result['text_color'] = run_info['font'].get('color')
            
            result['paragraphs'].append(para_info)
        
        return result
    
    def _get_font_color(self, font) -> str:
        """Get font color as hex string"""
        try:
            if font.color and font.color.rgb:
                return self._rgb_to_hex(font.color.rgb)
            elif font.color and font.color.theme_color:
                # Theme color - return a reasonable default based on theme
                return '#F8FAFC'
        except:
            pass
        return '#F8FAFC'
    
    def _get_chart_type(self, shape) -> str:
        """Extract chart type"""
        try:
            chart = shape.chart
            return str(chart.chart_type)
        except:
            return 'unknown'
    
    def _get_table_style(self, shape) -> Dict:
        """Extract table style"""
        try:
            table = shape.table
            style_info = {
                'rows': len(table.rows),
                'columns': len(table.columns),
                'first_row_header': True
            }
            
            # Try to get header row colors
            if len(table.rows) > 0:
                first_cell = table.cell(0, 0)
                if hasattr(first_cell, 'fill') and first_cell.fill.type:
                    fill_info = self._extract_fill_info(first_cell.fill)
                    if fill_info:
                        style_info['header_fill'] = fill_info
            
            return style_info
        except:
            return {}
    
    def _extract_theme_colors(self, prs: Presentation):
        """Extract theme colors from presentation"""
        try:
            theme_colors = []
            
            # Try to access theme
            if prs.slide_masters:
                master = prs.slide_masters[0]
                # Theme colors are typically in the master's theme
                # We extract from shapes as a fallback
                for shape in master.shapes:
                    if hasattr(shape, 'fill') and shape.fill.type:
                        try:
                            if shape.fill.fore_color and shape.fill.fore_color.rgb:
                                color = self._rgb_to_hex(shape.fill.fore_color.rgb)
                                if color not in theme_colors:
                                    theme_colors.append(color)
                        except:
                            pass
            
            if theme_colors:
                self.template_data['color_scheme'] = theme_colors[:10]
                
        except Exception as e:
            print(f"Error extracting theme colors: {e}")
    
    def _process_colors(self, colors: List[str]):
        """Process extracted colors"""
        if not colors:
            return
        
        valid_colors = [c for c in colors if c and re.match(r'^#[0-9A-Fa-f]{6}$', c)]
        
        if not valid_colors:
            return
        
        color_counts = Counter(valid_colors)
        sorted_colors = [color for color, _ in color_counts.most_common(10)]
        
        # Categorize by lightness
        dark_colors = []
        light_colors = []
        mid_colors = []
        
        for color in sorted_colors:
            lightness = self._get_lightness(color)
            if lightness < 0.3:
                dark_colors.append(color)
            elif lightness > 0.7:
                light_colors.append(color)
            else:
                mid_colors.append(color)
        
        # Assign colors
        if dark_colors:
            self.template_data['colors']['background'] = dark_colors[0]
        if light_colors:
            self.template_data['colors']['text_primary'] = light_colors[0]
            if len(light_colors) > 1:
                self.template_data['colors']['text_secondary'] = light_colors[1]
        if mid_colors:
            self.template_data['colors']['primary'] = mid_colors[0]
            if len(mid_colors) > 1:
                self.template_data['colors']['secondary'] = mid_colors[1]
            if len(mid_colors) > 2:
                self.template_data['colors']['accent'] = mid_colors[2]
    
    def _process_fonts(self, fonts: List[Dict]):
        """Process extracted fonts"""
        if not fonts:
            return
        
        # Group by size
        sized_fonts = [(f, f.get('size', 0)) for f in fonts if f.get('size')]
        sized_fonts.sort(key=lambda x: x[1], reverse=True)
        
        if sized_fonts:
            # Largest font = title
            title_font = sized_fonts[0][0]
            self.template_data['fonts']['title'] = {
                'name': title_font.get('name', 'Arial'),
                'size': title_font.get('size', 44),
                'bold': title_font.get('bold', True),
                'color': title_font.get('color', '#F8FAFC')
            }
            
            # Medium fonts = subtitle/body
            if len(sized_fonts) > 1:
                subtitle_font = sized_fonts[len(sized_fonts)//3][0]
                self.template_data['fonts']['subtitle'] = {
                    'name': subtitle_font.get('name', 'Arial'),
                    'size': subtitle_font.get('size', 28),
                    'bold': subtitle_font.get('bold', False),
                    'color': subtitle_font.get('color', '#94A3B8')
                }
            
            if len(sized_fonts) > 2:
                body_font = sized_fonts[-1][0]
                self.template_data['fonts']['body'] = {
                    'name': body_font.get('name', 'Arial'),
                    'size': body_font.get('size', 18),
                    'bold': body_font.get('bold', False),
                    'color': body_font.get('color', '#F8FAFC')
                }
    
    def analyze_image(self, image_file: io.BytesIO) -> Dict[str, Any]:
        """Analyze an image to extract styling"""
        try:
            image = Image.open(image_file)
            colors = self._extract_colors_from_image(image)
            self._process_colors(colors)
            return self.template_data
        except Exception as e:
            print(f"Error analyzing image: {e}")
            return self.template_data
    
    def _extract_colors_from_image(self, image: Image.Image) -> List[str]:
        """Extract dominant colors from image"""
        image = image.resize((150, 150))
        
        if image.mode != 'RGB':
            image = image.convert('RGB')
        
        pixels = list(image.getdata())
        color_counts = Counter()
        
        for pixel in pixels:
            rounded = (pixel[0] // 32 * 32, pixel[1] // 32 * 32, pixel[2] // 32 * 32)
            color_counts[rounded] += 1
        
        top_colors = color_counts.most_common(10)
        return [self._rgb_tuple_to_hex(color) for color, _ in top_colors]
    
    def _rgb_to_hex(self, rgb: RGBColor) -> str:
        """Convert RGBColor to hex"""
        try:
            return f'#{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}'
        except:
            return '#000000'
    
    def _rgb_tuple_to_hex(self, rgb: Tuple[int, int, int]) -> str:
        """Convert RGB tuple to hex"""
        return f'#{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}'
    
    def _get_lightness(self, hex_color: str) -> float:
        """Get lightness of hex color"""
        hex_color = hex_color.lstrip('#')
        r = int(hex_color[0:2], 16) / 255
        g = int(hex_color[2:4], 16) / 255
        b = int(hex_color[4:6], 16) / 255
        _, lightness, _ = colorsys.rgb_to_hls(r, g, b)
        return lightness
    
    def get_template_summary(self) -> str:
        """Generate summary of analyzed template"""
        summary = ["## Template Analysis Summary\n"]
        
        summary.append("### Colors")
        for name, color in self.template_data['colors'].items():
            summary.append(f"- {name}: {color}")
        
        summary.append("\n### Fonts")
        for name, font in self.template_data['fonts'].items():
            summary.append(f"- {name}: {font.get('name', 'Arial')} {font.get('size', 18)}pt")
        
        summary.append(f"\n### Logo: {'Yes ✓' if self.template_data['has_logo'] else 'No'}")
        summary.append(f"### Background: {self.template_data['background'].get('type', 'solid')}")
        
        return "\n".join(summary)
