"""
Enhanced PowerPoint Generator with complete layout support
"""
from typing import Dict, List, Any, Optional
from pathlib import Path
import io
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RgbColor
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image
import base64
import logging

logger = logging.getLogger(__name__)

class EnhancedPPTGenerator:
    """Professional PowerPoint generator with full layout support"""
    
    # Layout mappings - adjust based on your template
    LAYOUT_MAP = {
        'title': 0,           # Title Slide
        'content': 1,         # Title and Content
        'two_column': 1,      # Use Title and Content (we'll create columns manually)
        'chart': 1,           # Title and Content
        'table': 1,           # Title and Content
        'quote': 1,           # Title and Content
        'metrics': 1,         # Title and Content
        'timeline': 1,        # Title and Content
        'comparison': 1,      # Title and Content
        'conclusion': 1,      # Title and Content
        'image': 1,           # Title and Content
    }
    
    def __init__(self, template_path: Optional[str] = None):
        """Initialize with optional template"""
        if template_path and Path(template_path).exists():
            self.prs = Presentation(template_path)
        else:
            self.prs = Presentation()
        
        self.slide_width = self.prs.slide_width
        self.slide_height = self.prs.slide_height
    
    def generate(self, slides_data: List[Dict[str, Any]]) -> io.BytesIO:
        """
        Generate complete presentation from slides data
        
        Args:
            slides_data: List of slide dictionaries with layout, title, content
        
        Returns:
            BytesIO object with PPTX file
        """
        generated_count = 0
        failed_slides = []
        
        for slide_data in slides_data:
            try:
                self._create_slide(slide_data)
                generated_count += 1
            except Exception as e:
                slide_num = slide_data.get('slide_number', 'Unknown')
                layout = slide_data.get('layout', 'Unknown')
                error_msg = f"Failed to generate slide {slide_num} ({layout}): {str(e)}"
                logger.error(error_msg, exc_info=True)
                failed_slides.append({'slide': slide_num, 'layout': layout, 'error': str(e)})
                
                # Create error placeholder slide
                self._create_error_slide(slide_num, layout, str(e))
        
        if failed_slides:
            logger.warning(f"Generated {generated_count} slides. Failed: {len(failed_slides)}")
            logger.warning(f"Failed slides: {failed_slides}")
        
        # Save to BytesIO
        output = io.BytesIO()
        self.prs.save(output)
        output.seek(0)
        
        return output
    
    def _create_slide(self, slide_data: Dict[str, Any]):
        """Create a single slide based on layout type"""
        layout_type = slide_data.get('layout', 'content')
        layout_idx = self.LAYOUT_MAP.get(layout_type, 1)
        
        # Ensure layout index is valid
        if layout_idx >= len(self.prs.slide_layouts):
            layout_idx = 1  # Fallback to default
        
        slide_layout = self.prs.slide_layouts[layout_idx]
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Set title if exists
        if slide_data.get('title') and hasattr(slide, 'shapes') and slide.shapes.title:
            slide.shapes.title.text = slide_data['title']
        
        # Route to specific layout handler
        handler = getattr(self, f'_handle_{layout_type}_layout', self._handle_content_layout)
        handler(slide, slide_data)
    
    def _handle_title_layout(self, slide, slide_data):
        """Handle title slide layout"""
        content = slide_data.get('content', {})
        
        if slide.shapes.title:
            slide.shapes.title.text = slide_data.get('title', '')
        
        # Add subtitle if exists
        if len(slide.placeholders) > 1:
            subtitle = slide.placeholders[1]
            subtitle.text = slide_data.get('subtitle', content.get('subtitle', ''))
    
    def _handle_content_layout(self, slide, slide_data):
        """Handle standard content layout with bullet points"""
        content = slide_data.get('content', {})
        
        # Add main text or bullet points
        if hasattr(slide, 'shapes') and len(slide.shapes) > 1:
            body_shape = slide.shapes[1]  # Usually the content placeholder
            
            if body_shape.has_text_frame:
                tf = body_shape.text_frame
                tf.clear()
                
                # Add main text
                if content.get('main_text'):
                    p = tf.paragraphs[0]
                    p.text = content['main_text']
                    p.font.size = Pt(18)
                
                # Add bullet points
                for bullet in content.get('bullet_points', []):
                    p = tf.add_paragraph()
                    p.text = f"• {bullet}"
                    p.font.size = Pt(16)
                    p.level = 0
                
                # Add speaker notes if exists
                if slide_data.get('speaker_notes'):
                    notes_slide = slide.notes_slide
                    notes_slide.notes_text_frame.text = slide_data['speaker_notes']
    
    def _handle_two_column_layout(self, slide, slide_data):
        """Handle two-column layout - CRITICAL FIX"""
        content = slide_data.get('content', {})
        
        # Remove default content placeholder if exists
        if len(slide.shapes) > 1:
            sp = slide.shapes[1]
            slide.shapes._spTree.remove(sp._element)
        
        # Calculate column dimensions
        left_margin = Inches(0.5)
        right_margin = Inches(0.5)
        top_pos = Inches(2.0)
        bottom_margin = Inches(0.5)
        available_width = self.slide_width - left_margin - right_margin
        col_width = (available_width - Inches(0.3)) / 2  # Gap between columns
        
        # Left column
        left_col = slide.shapes.add_textbox(
            left_margin, top_pos, col_width, 
            self.slide_height - top_pos - bottom_margin
        )
        tf_left = left_col.text_frame
        tf_left.word_wrap = True
        
        # Add left column content
        left_content = content.get('left_column', '')
        if isinstance(left_content, list):
            for i, item in enumerate(left_content):
                p = tf_left.paragraphs[0] if i == 0 else tf_left.add_paragraph()
                p.text = f"• {item}" if i > 0 or not item.startswith('•') else item
                p.font.size = Pt(14)
        else:
            tf_left.text = left_content
            tf_left.paragraphs[0].font.size = Pt(14)
        
        # Right column
        right_col = slide.shapes.add_textbox(
            left_margin + col_width + Inches(0.3), top_pos, col_width,
            self.slide_height - top_pos - bottom_margin
        )
        tf_right = right_col.text_frame
        tf_right.word_wrap = True
        
        # Add right column content
        right_content = content.get('right_column', '')
        if isinstance(right_content, list):
            for i, item in enumerate(right_content):
                p = tf_right.paragraphs[0] if i == 0 else tf_right.add_paragraph()
                p.text = f"• {item}" if i > 0 or not item.startswith('•') else item
                p.font.size = Pt(14)
        else:
            tf_right.text = right_content
            tf_right.paragraphs[0].font.size = Pt(14)
    
    def _handle_chart_layout(self, slide, slide_data):
        """Handle chart layout - CRITICAL FIX"""
        content = slide_data.get('content', {})
        chart_data = content.get('chart', {})
        
        if not chart_data:
            # Fallback to text if no chart data
            self._handle_content_layout(slide, slide_data)
            return
        
        # Remove default content placeholder
        if len(slide.shapes) > 1:
            sp = slide.shapes[1]
            slide.shapes._spTree.remove(sp._element)
        
        # Create chart placeholder box
        chart_width = self.slide_width * 0.8
        chart_height = self.slide_height * 0.55
        left_pos = (self.slide_width - chart_width) / 2
        top_pos = Inches(2.0)
        
        # Add chart shape (placeholder)
        chart_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            left_pos, top_pos, chart_width, chart_height
        )
        chart_shape.fill.solid()
        chart_shape.fill.fore_color.rgb = RgbColor(240, 240, 240)
        chart_shape.line.color.rgb = RgbColor(200, 200, 200)
        
        # Add chart title inside shape
        tf = chart_shape.text_frame
        tf.text = f"📊 {chart_data.get('title', 'Chart')}"
        tf.paragraphs[0].font.size = Pt(16)
        tf.paragraphs[0].font.bold = True
        
        # Add chart description
        if chart_data.get('description'):
            p = tf.add_paragraph()
            p.text = f"\n{chart_data['description']}"
            p.font.size = Pt(12)
        
        # Add chart data as text below
        chart_info = slide.shapes.add_textbox(
            Inches(0.5), top_pos + chart_height + Inches(0.2),
            self.slide_width - Inches(1), Inches(1.5)
        )
        tf_info = chart_info.text_frame
        
        if 'data' in chart_data:
            p = tf_info.paragraphs[0]
            p.text = "Data Points:"
            p.font.bold = True
            p.font.size = Pt(12)
            
            for key, value in chart_data['data'].items():
                p = tf_info.add_paragraph()
                p.text = f"  • {key}: {value}"
                p.font.size = Pt(10)
    
    def _handle_table_layout(self, slide, slide_data):
        """Handle table layout"""
        content = slide_data.get('content', {})
        table_data = content.get('table', {})
        
        if not table_data:
            self._handle_content_layout(slide, slide_data)
            return
        
        # Remove default content placeholder
        if len(slide.shapes) > 1:
            sp = slide.shapes[1]
            slide.shapes._spTree.remove(sp._element)
        
        # Create table
        rows = len(table_data.get('data', [])) + 1  # +1 for header
        cols = len(table_data.get('headers', []))
        
        if rows == 1 or cols == 0:
            logger.warning("Invalid table data")
            return
        
        table_width = self.slide_width * 0.9
        table_height = self.slide_height * 0.6
        left_pos = (self.slide_width - table_width) / 2
        top_pos = Inches(1.8)
        
        table = slide.shapes.add_table(
            rows, cols, left_pos, top_pos, table_width, table_height
        ).table
        
        # Set column widths
        col_width = table_width / cols
        for i in range(cols):
            table.columns[i].width = col_width
        
        # Add headers
        for i, header in enumerate(table_data.get('headers', [])):
            cell = table.cell(0, i)
            cell.text = header
            cell.text_frame.paragraphs[0].font.bold = True
            cell.text_frame.paragraphs[0].font.size = Pt(12)
        
        # Add data rows
        for row_idx, row_data in enumerate(table_data.get('data', []), 1):
            for col_idx, value in enumerate(row_data):
                if col_idx < cols:
                    cell = table.cell(row_idx, col_idx)
                    cell.text = str(value)
                    cell.text_frame.paragraphs[0].font.size = Pt(10)
    
    def _handle_quote_layout(self, slide, slide_data):
        """Handle quote layout"""
        content = slide_data.get('content', {})
        
        # Remove default content placeholder
        if len(slide.shapes) > 1:
            sp = slide.shapes[1]
            slide.shapes._spTree.remove(sp._element)
        
        # Add quote shape
        quote_box = slide.shapes.add_shape(
            MSO_SHAPE.TEXT_BOX,
            Inches(0.8), Inches(2.2),
            self.slide_width - Inches(1.6), Inches(2.5)
        )
        quote_box.fill.solid()
        quote_box.fill.fore_color.rgb = RgbColor(245, 245, 245)
        quote_box.line.color.rgb = RgbColor(200, 200, 200)
        
        tf = quote_box.text_frame
        tf.word_wrap = True
        
        # Add opening quote
        p = tf.paragraphs[0]
        p.text = f""{content.get('quote', 'Quote text')}""
        p.font.size = Pt(20)
        p.font.italic = True
        
        # Add author
        if content.get('quote_author'):
            p_author = tf.add_paragraph()
            p_author.text = f"\n— {content['quote_author']}"
            p_author.font.size = Pt(14)
            p_author.alignment = PP_ALIGN.RIGHT
    
    def _handle_metrics_layout(self, slide, slide_data):
        """Handle KPI/metrics layout"""
        content = slide_data.get('content', {})
        metrics = content.get('key_metrics', [])
        
        if not metrics:
            self._handle_content_layout(slide, slide_data)
            return
        
        # Remove default content placeholder
        if len(slide.shapes) > 1:
            sp = slide.shapes[1]
            slide.shapes._spTree.remove(sp._element)
        
        # Create metric boxes
        num_metrics = len(metrics)
        box_width = (self.slide_width - Inches(1)) / num_metrics - Inches(0.2)
        box_height = Inches(2.5)
        top_pos = Inches(2.2)
        left_margin = Inches(0.5)
        
        for i, metric in enumerate(metrics):
            left_pos = left_margin + (i * (box_width + Inches(0.2)))
            
            # Add metric box
            metric_box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                left_pos, top_pos, box_width, box_height
            )
            metric_box.fill.solid()
            metric_box.fill.fore_color.rgb = RgbColor(79, 129, 189)
            metric_box.line.color.rgb = RgbColor(79, 129, 189)
            
            # Add metric value
            tf = metric_box.text_frame
            tf.clear()
            
            p_value = tf.paragraphs[0]
            p_value.text = metric.get('value', 'N/A')
            p_value.font.size = Pt(28)
            p_value.font.bold = True
            p_value.alignment = PP_ALIGN.CENTER
            
            p_label = tf.add_paragraph()
            p_label.text = metric.get('label', '')
            p_label.font.size = Pt(14)
            p_label.alignment = PP_ALIGN.CENTER
    
    def _handle_timeline_layout(self, slide, slide_data):
        """Handle timeline layout"""
        content = slide_data.get('content', {})
        timeline_items = content.get('timeline_items', [])
        
        if not timeline_items:
            self._handle_content_layout(slide, slide_data)
            return
        
        # Remove default content placeholder
        if len(slide.shapes) > 1:
            sp = slide.shapes[1]
            slide.shapes._spTree.remove(sp._element)
        
        # Create timeline
        timeline_width = self.slide_width * 0.9
        timeline_left = (self.slide_width - timeline_width) / 2
        timeline_y = Inches(3.0)
        
        # Draw timeline line
        line = slide.shapes.add_shape(
            MSO_SHAPE.LINE,
            timeline_left, timeline_y + Inches(0.5),
            timeline_width, Inches(0.05)
        )
        line.line.color.rgb = RgbColor(79, 129, 189)
        line.line.width = Pt(3)
        
        # Add timeline items
        num_items = len(timeline_items)
        spacing = timeline_width / (num_items if num_items > 1 else 1)
        
        for i, item in enumerate(timeline_items):
            x_pos = timeline_left + (i * spacing) + (spacing / 2)
            
            # Add circle marker
            marker = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                x_pos - Inches(0.25), timeline_y + Inches(0.3),
                Inches(0.5), Inches(0.5)
            )
            marker.fill.solid()
            marker.fill.fore_color.rgb = RgbColor(79, 129, 189)
            
            # Add date/title
            tf_box = slide.shapes.add_textbox(
                x_pos - Inches(1), timeline_y - Inches(0.8),
                Inches(2), Inches(0.6)
            )
            tf = tf_box.text_frame
            p = tf.paragraphs[0]
            p.text = item.get('date', '')
            p.font.bold = True
            p.font.size = Pt(12)
            p.alignment = PP_ALIGN.CENTER
            
            # Add description
            if item.get('description'):
                desc_box = slide.shapes.add_textbox(
                    x_pos - Inches(1.5), timeline_y + Inches(1.0),
                    Inches(3), Inches(1.5)
                )
                tf_desc = desc_box.text_frame
                tf_desc.word_wrap = True
                p_desc = tf_desc.paragraphs[0]
                p_desc.text = item['description']
                p_desc.font.size = Pt(10)
                p_desc.alignment = PP_ALIGN.CENTER
    
    def _handle_comparison_layout(self, slide, slide_data):
        """Handle comparison layout"""
        content = slide_data.get('content', {})
        
        # Use two column handler as base
        self._handle_two_column_layout(slide, slide_data)
    
    def _handle_image_layout(self, slide, slide_data):
        """Handle image layout - CRITICAL FIX"""
        content = slide_data.get('content', {})
        image_data = content.get('image', '')
        
        # Remove default content placeholder
        if len(slide.shapes) > 1:
            sp = slide.shapes[1]
            slide.shapes._spTree.remove(sp._element)
        
        # Calculate image dimensions
        img_width = self.slide_width * 0.7
        img_height = self.slide_height * 0.55
        left_pos = (self.slide_width - img_width) / 2
        top_pos = Inches(1.8)
        
        try:
            # Handle different image input formats
            if isinstance(image_data, str):
                if image_data.startswith('data:image'):
                    # Base64 encoded image
                    import re
                    match = re.search(r'data:image/\w+;base64,(.+)', image_data)
                    if match:
                        img_bytes = base64.b64decode(match.group(1))
                        img_stream = io.BytesIO(img_bytes)
                        slide.shapes.add_picture(img_stream, left_pos, top_pos, 
                                                width=img_width, height=img_height)
                    else:
                        raise ValueError("Invalid base64 image data")
                elif Path(image_data).exists():
                    # File path
                    slide.shapes.add_picture(image_data, left_pos, top_pos,
                                            width=img_width, height=img_height)
                else:
                    raise ValueError(f"Image not found: {image_data}")
            elif isinstance(image_data, bytes):
                # Raw bytes
                img_stream = io.BytesIO(image_data)
                slide.shapes.add_picture(img_stream, left_pos, top_pos,
                                        width=img_width, height=img_height)
            else:
                raise ValueError(f"Unsupported image format: {type(image_data)}")
        except Exception as e:
            logger.error(f"Failed to add image: {e}", exc_info=True)
            # Add placeholder instead
            placeholder = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                left_pos, top_pos, img_width, img_height
            )
            placeholder.fill.solid()
            placeholder.fill.fore_color.rgb = RgbColor(240, 240, 240)
            tf = placeholder.text_frame
            tf.text = "🖼️ Image Placeholder\n(Image failed to load)"
            tf.paragraphs[0].font.size = Pt(16)
    
    def _handle_conclusion_layout(self, slide, slide_data):
        """Handle conclusion slide"""
        content = slide_data.get('content', {})
        
        # Remove default content placeholder
        if len(slide.shapes) > 1:
            sp = slide.shapes[1]
            slide.shapes._spTree.remove(sp._element)
        
        # Add conclusion box
        conclusion_box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.5), Inches(2.0),
            self.slide_width - Inches(1), Inches(3.0)
        )
        conclusion_box.fill.solid()
        conclusion_box.fill.fore_color.rgb = RgbColor(79, 129, 189)
        
        tf = conclusion_box.text_frame
        tf.word_wrap = True
        
        # Add main conclusion text
        main_text = content.get('main_text', '')
        if main_text:
            p = tf.paragraphs[0]
            p.text = main_text
            p.font.size = Pt(24)
            p.font.bold = True
            p.alignment = PP_ALIGN.CENTER
        
        # Add call to action or key points
        for bullet in content.get('bullet_points', []):
            p = tf.add_paragraph()
            p.text = f"\n✓ {bullet}"
            p.font.size = Pt(16)
            p.alignment = PP_ALIGN.CENTER
    
    def _create_error_slide(self, slide_num: int, layout: str, error: str):
        """Create error placeholder slide when generation fails"""
        layout_idx = 1 if len(self.prs.slide_layouts) > 1 else 0
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[layout_idx])
        
        if slide.shapes.title:
            slide.shapes.title.text = f"⚠️ Slide {slide_num} - Generation Failed"
        
        # Add error message
        if len(slide.shapes) > 1:
            error_box = slide.shapes[1]
            tf = error_box.text_frame
            tf.text = f"Layout: {layout}\n\nError: {error}\n\nPlease check logs for details."
            tf.paragraphs[0].font.size = Pt(14)
            tf.paragraphs[0].font.color.rgb = RgbColor(255, 0, 0)
