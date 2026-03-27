"""
Template Cloner - Complete Content Generation
Handles ALL layout types with full content
"""

import io
from typing import Dict, List, Optional
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

from .chart_generator import ChartGenerator


class TemplateCloner:
    """Clones template and generates complete slide content"""
    
    def __init__(self, template_bytes: Optional[bytes] = None):
        self.template_bytes = template_bytes
        self.prs = None
        self.available_layouts = []
        self.chart_generator = ChartGenerator()
        
        if template_bytes:
            self._load_template()
    
    def _load_template(self):
        """Load the template presentation"""
        try:
            template_stream = io.BytesIO(self.template_bytes)
            self.prs = Presentation(template_stream)
            
            self.available_layouts = []
            for idx, layout in enumerate(self.prs.slide_layouts):
                self.available_layouts.append({
                    'index': idx,
                    'name': layout.name,
                    'layout': layout
                })
        except Exception as e:
            print(f"Error loading template: {e}")
            self.prs = Presentation()
    
    def generate_presentation(self, content: Dict) -> io.BytesIO:
        """Generate complete presentation"""
        
        if not self.prs:
            self._load_template()
        
        if not self.prs:
            self.prs = Presentation()
        
        # Clear existing slides
        self._clear_content_slides()
        
        # Generate slides
        slides_data = content.get('slides', [])
        
        for slide_data in slides_data:
            self._add_slide_from_content(slide_data)
        
        # Save
        output = io.BytesIO()
        self.prs.save(output)
        output.seek(0)
        
        return output
    
    def _clear_content_slides(self):
        """Remove all slides"""
        try:
            slide_count = len(self.prs.slides)
            
            for _ in range(slide_count):
                if len(self.prs.slides) > 0:
                    xml_slides = self.prs.slides._sldIdLst
                    if len(xml_slides) > 0:
                        slide_id = xml_slides[0]
                        try:
                            self.prs.part.drop_rel(slide_id.rId)
                        except:
                            pass
                        xml_slides.remove(slide_id)
        except Exception as e:
            print(f"Warning: Could not clear slides: {e}")
    
    def _add_slide_from_content(self, slide_data: Dict):
        """Add slide with complete content"""
        
        layout_type = slide_data.get('layout', 'content')
        layout = self._get_layout_for_type(layout_type)
        slide = self.prs.slides.add_slide(layout)
        
        # Fill content based on layout type
        handlers = {
            'title': self._fill_title_slide,
            'content': self._fill_content_slide,
            'two_column': self._fill_two_column_slide,
            'chart': self._fill_chart_slide,
            'table': self._fill_table_slide,
            'quote': self._fill_quote_slide,
            'metrics': self._fill_metrics_slide,
            'timeline': self._fill_timeline_slide,
            'comparison': self._fill_comparison_slide,
            'conclusion': self._fill_conclusion_slide,
            'image': self._fill_image_slide
        }
        
        handler = handlers.get(layout_type, self._fill_content_slide)
        handler(slide, slide_data)
    
    def _get_layout_for_type(self, layout_type: str):
        """Get matching layout"""
        
        layout_patterns = {
            'title': ['title slide', 'title', 'cover'],
            'content': ['title and content', 'content', 'bullet'],
            'two_column': ['two content', 'comparison', 'two column'],
            'chart': ['title and content', 'content', 'chart'],
            'table': ['title and content', 'content', 'table'],
            'quote': ['blank', 'quote', 'title only'],
            'conclusion': ['title slide', 'section', 'conclusion'],
            'metrics': ['title and content', 'content'],
            'timeline': ['title and content', 'content'],
            'comparison': ['comparison', 'two content'],
            'image': ['picture', 'content']
        }
        
        preferred = layout_patterns.get(layout_type, ['title and content'])
        
        for pattern in preferred:
            for layout_info in self.available_layouts:
                if pattern.lower() in layout_info['name'].lower():
                    return layout_info['layout']
        
        if self.available_layouts:
            if layout_type == 'title':
                return self.available_layouts[0]['layout']
            elif len(self.available_layouts) > 1:
                return self.available_layouts[1]['layout']
            else:
                return self.available_layouts[0]['layout']
        
        return self.prs.slide_layouts[0]
    
    def _fill_placeholders(self, slide, title: str = "", content_text: str = ""):
        """Fill standard placeholders"""
        for shape in slide.shapes:
            try:
                if not hasattr(shape, 'is_placeholder') or not shape.is_placeholder:
                    continue
                
                placeholder_type = str(shape.placeholder_format.type)
                
                if 'TITLE' in placeholder_type and title:
                    if shape.has_text_frame:
                        shape.text_frame.clear()
                        shape.text_frame.paragraphs[0].text = title
                
                elif 'BODY' in placeholder_type or 'OBJECT' in placeholder_type:
                    if shape.has_text_frame and content_text:
                        shape.text_frame.clear()
                        shape.text_frame.paragraphs[0].text = content_text
                        
            except Exception:
                continue
    
    def _add_title(self, slide, title: str, top: float = 0.4):
        """Add title text box"""
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(top),
            Inches(self.prs.slide_width.inches - 1),
            Inches(1)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(36)
        p.font.bold = True
    
    def _add_bullets(self, slide, bullets: List[str], left: float = 0.75, top: float = 1.5, width: float = None, height: float = 5):
        """Add bullet points"""
        if width is None:
            width = self.prs.slide_width.inches - 1.5
        
        text_box = slide.shapes.add_textbox(
            Inches(left), Inches(top),
            Inches(width), Inches(height)
        )
        tf = text_box.text_frame
        tf.word_wrap = True
        
        for i, bullet in enumerate(bullets):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            
            p.text = f"• {bullet}"
            p.font.size = Pt(18)
            p.space_after = Pt(10)
    
    def _fill_title_slide(self, slide, slide_data: Dict):
        """Fill title slide"""
        title = slide_data.get('title', '')
        subtitle = slide_data.get('subtitle', '')
        
        # Try placeholders first
        self._fill_placeholders(slide, title, subtitle)
        
        # Check if filled
        title_filled = False
        for shape in slide.shapes:
            if hasattr(shape, 'text_frame') and title in shape.text_frame.text:
                title_filled = True
                break
        
        if not title_filled:
            # Add manually
            self._add_title(slide, title, 2.5)
            
            if subtitle:
                sub_box = slide.shapes.add_textbox(
                    Inches(0.5), Inches(4),
                    Inches(self.prs.slide_width.inches - 1),
                    Inches(1)
                )
                tf = sub_box.text_frame
                p = tf.paragraphs[0]
                p.text = subtitle
                p.font.size = Pt(24)
                p.alignment = PP_ALIGN.CENTER
    
    def _fill_content_slide(self, slide, slide_data: Dict):
        """Fill content slide with bullets"""
        title = slide_data.get('title', '')
        content_data = slide_data.get('content', {})
        bullets = content_data.get('bullet_points', [])
        main_text = content_data.get('main_text', '')
        
        # Fill placeholders
        content_text = '\n'.join([f"• {b}" for b in bullets]) if bullets else main_text
        self._fill_placeholders(slide, title, content_text)
        
        # Check if content filled
        content_filled = False
        for shape in slide.shapes:
            if hasattr(shape, 'text_frame'):
                text = shape.text_frame.text
                if bullets and bullets[0][:20] in text:
                    content_filled = True
                    break
        
        if not content_filled:
            self._add_title(slide, title)
            if bullets:
                self._add_bullets(slide, bullets)
            elif main_text:
                text_box = slide.shapes.add_textbox(
                    Inches(0.75), Inches(1.5),
                    Inches(self.prs.slide_width.inches - 1.5),
                    Inches(5)
                )
                tf = text_box.text_frame
                tf.word_wrap = True
                tf.paragraphs[0].text = main_text
                tf.paragraphs[0].font.size = Pt(18)
    
    def _fill_two_column_slide(self, slide, slide_data: Dict):
        """Fill two-column slide"""
        title = slide_data.get('title', '')
        content_data = slide_data.get('content', {})
        
        left_content = content_data.get('left_column', [])
        right_content = content_data.get('right_column', [])
        
        # Add title
        self._add_title(slide, title)
        
        slide_width = self.prs.slide_width.inches
        col_width = (slide_width - 1.5) / 2
        
        # Left column
        if isinstance(left_content, list):
            self._add_bullets(slide, left_content, 0.5, 1.5, col_width)
        elif left_content:
            text_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(1.5),
                Inches(col_width), Inches(5)
            )
            tf = text_box.text_frame
            tf.word_wrap = True
            tf.paragraphs[0].text = str(left_content)
            tf.paragraphs[0].font.size = Pt(16)
        
        # Right column
        if isinstance(right_content, list):
            self._add_bullets(slide, right_content, slide_width/2 + 0.25, 1.5, col_width)
        elif right_content:
            text_box = slide.shapes.add_textbox(
                Inches(slide_width/2 + 0.25), Inches(1.5),
                Inches(col_width), Inches(5)
            )
            tf = text_box.text_frame
            tf.word_wrap = True
            tf.paragraphs[0].text = str(right_content)
            tf.paragraphs[0].font.size = Pt(16)
        
        # Divider line
        try:
            line = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(slide_width/2 - 0.02), Inches(1.5),
                Inches(0.04), Inches(4.5)
            )
            line.fill.solid()
            line.fill.fore_color.rgb = RGBColor(100, 100, 100)
            line.line.fill.background()
        except:
            pass
    
    def _fill_chart_slide(self, slide, slide_data: Dict):
        """Fill chart slide"""
        title = slide_data.get('title', '')
        content_data = slide_data.get('content', {})
        chart_info = content_data.get('chart', {})
        
        # Add title
        self._add_title(slide, title)
        
        # Generate chart
        if chart_info:
            chart_type = chart_info.get('type', 'bar')
            chart_data = {
                'title': chart_info.get('title', ''),
                'labels': chart_info.get('labels', ['A', 'B', 'C', 'D']),
                'datasets': chart_info.get('datasets', [{'name': 'Series', 'values': [65, 78, 90, 85]}]),
                'x_axis_label': chart_info.get('x_axis_label', ''),
                'y_axis_label': chart_info.get('y_axis_label', '')
            }
            
            try:
                chart_image = self.chart_generator.create_chart(chart_type, chart_data, 800, 400)
                image_stream = io.BytesIO(chart_image)
                
                slide.shapes.add_picture(
                    image_stream,
                    Inches(1.5), Inches(1.8),
                    Inches(self.prs.slide_width.inches - 3),
                    Inches(4.5)
                )
            except Exception as e:
                print(f"Chart error: {e}")
                # Add placeholder
                text_box = slide.shapes.add_textbox(
                    Inches(2), Inches(3),
                    Inches(self.prs.slide_width.inches - 4),
                    Inches(1)
                )
                tf = text_box.text_frame
                tf.paragraphs[0].text = f"[Chart: {chart_info.get('title', 'Data Visualization')}]"
                tf.paragraphs[0].font.size = Pt(24)
                tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    def _fill_table_slide(self, slide, slide_data: Dict):
        """Fill table slide"""
        title = slide_data.get('title', '')
        content_data = slide_data.get('content', {})
        table_info = content_data.get('table', {})
        
        # Add title
        self._add_title(slide, title)
        
        headers = table_info.get('headers', [])
        rows = table_info.get('rows', [])
        
        if headers:
            num_cols = len(headers)
            num_rows = len(rows) + 1
            
            table_width = min(self.prs.slide_width.inches - 2, num_cols * 2)
            
            table = slide.shapes.add_table(
                num_rows, num_cols,
                Inches(1), Inches(1.8),
                Inches(table_width),
                Inches(min(num_rows * 0.5, 4.5))
            ).table
            
            # Header row
            for i, header in enumerate(headers):
                cell = table.cell(0, i)
                cell.text = str(header)
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(59, 130, 246)
                
                para = cell.text_frame.paragraphs[0]
                para.font.bold = True
                para.font.color.rgb = RGBColor(255, 255, 255)
                para.font.size = Pt(12)
                para.alignment = PP_ALIGN.CENTER
            
            # Data rows
            for row_idx, row_data in enumerate(rows):
                for col_idx, cell_data in enumerate(row_data):
                    if col_idx < num_cols:
                        cell = table.cell(row_idx + 1, col_idx)
                        cell.text = str(cell_data)
                        
                        para = cell.text_frame.paragraphs[0]
                        para.font.size = Pt(11)
                        para.alignment = PP_ALIGN.CENTER
    
    def _fill_quote_slide(self, slide, slide_data: Dict):
        """Fill quote slide"""
        content_data = slide_data.get('content', {})
        quote = content_data.get('quote', '')
        author = content_data.get('quote_author', '')
        
        slide_width = self.prs.slide_width.inches
        
        # Quote text
        quote_box = slide.shapes.add_textbox(
            Inches(1.5), Inches(2.5),
            Inches(slide_width - 3),
            Inches(3)
        )
        tf = quote_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f'"{quote}"'
        p.font.size = Pt(28)
        p.font.italic = True
        p.alignment = PP_ALIGN.CENTER
        
        # Author
        if author:
            author_box = slide.shapes.add_textbox(
                Inches(1.5), Inches(5.5),
                Inches(slide_width - 3),
                Inches(0.5)
            )
            tf = author_box.text_frame
            p = tf.paragraphs[0]
            p.text = f"— {author}"
            p.font.size = Pt(20)
            p.alignment = PP_ALIGN.RIGHT
    
    def _fill_metrics_slide(self, slide, slide_data: Dict):
        """Fill metrics/KPI slide"""
        title = slide_data.get('title', 'Key Metrics')
        content_data = slide_data.get('content', {})
        metrics = content_data.get('key_metrics', [])
        
        # Add title
        self._add_title(slide, title)
        
        if metrics:
            num_metrics = min(len(metrics), 4)
            slide_width = self.prs.slide_width.inches
            metric_width = (slide_width - 3) / num_metrics
            
            for i, metric in enumerate(metrics[:4]):
                x_pos = 1.5 + (i * metric_width)
                
                # Box
                try:
                    box = slide.shapes.add_shape(
                        MSO_SHAPE.ROUNDED_RECTANGLE,
                        Inches(x_pos), Inches(2.5),
                        Inches(metric_width - 0.3), Inches(3)
                    )
                    box.fill.solid()
                    box.fill.fore_color.rgb = RGBColor(30, 41, 59)
                    box.line.color.rgb = RGBColor(99, 102, 241)
                except:
                    pass
                
                # Value
                value_box = slide.shapes.add_textbox(
                    Inches(x_pos), Inches(3),
                    Inches(metric_width - 0.3), Inches(1)
                )
                tf = value_box.text_frame
                p = tf.paragraphs[0]
                p.text = str(metric.get('value', ''))
                p.font.size = Pt(32)
                p.font.bold = True
                p.font.color.rgb = RGBColor(236, 72, 153)
                p.alignment = PP_ALIGN.CENTER
                
                # Label
                label_box = slide.shapes.add_textbox(
                    Inches(x_pos), Inches(4.2),
                    Inches(metric_width - 0.3), Inches(0.8)
                )
                tf = label_box.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = metric.get('label', '')
                p.font.size = Pt(14)
                p.alignment = PP_ALIGN.CENTER
    
    def _fill_timeline_slide(self, slide, slide_data: Dict):
        """Fill timeline slide"""
        title = slide_data.get('title', 'Timeline')
        content_data = slide_data.get('content', {})
        items = content_data.get('timeline_items', [])
        
        # Add title
        self._add_title(slide, title)
        
        if items:
            slide_width = self.prs.slide_width.inches
            
            # Timeline line
            try:
                line = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    Inches(1), Inches(4),
                    Inches(slide_width - 2), Inches(0.05)
                )
                line.fill.solid()
                line.fill.fore_color.rgb = RGBColor(99, 102, 241)
                line.line.fill.background()
            except:
                pass
            
            # Items
            item_count = min(len(items), 5)
            item_width = (slide_width - 2) / item_count
            
            for i, item in enumerate(items[:5]):
                x_pos = 1 + (i * item_width)
                
                # Circle
                try:
                    circle = slide.shapes.add_shape(
                        MSO_SHAPE.OVAL,
                        Inches(x_pos + item_width/2 - 0.15), Inches(3.85),
                        Inches(0.3), Inches(0.3)
                    )
                    circle.fill.solid()
                    circle.fill.fore_color.rgb = RGBColor(236, 72, 153)
                    circle.line.fill.background()
                except:
                    pass
                
                # Date
                date_box = slide.shapes.add_textbox(
                    Inches(x_pos), Inches(2.5),
                    Inches(item_width), Inches(0.5)
                )
                tf = date_box.text_frame
                p = tf.paragraphs[0]
                p.text = item.get('date', '')
                p.font.size = Pt(12)
                p.font.bold = True
                p.alignment = PP_ALIGN.CENTER
                
                # Event
                event_box = slide.shapes.add_textbox(
                    Inches(x_pos), Inches(4.3),
                    Inches(item_width), Inches(1.5)
                )
                tf = event_box.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = item.get('event', '')
                p.font.size = Pt(10)
                p.alignment = PP_ALIGN.CENTER
    
    def _fill_comparison_slide(self, slide, slide_data: Dict):
        """Fill comparison slide"""
        title = slide_data.get('title', 'Comparison')
        content_data = slide_data.get('content', {})
        items = content_data.get('comparison_items', [])
        
        # Add title
        self._add_title(slide, title)
        
        if len(items) >= 2:
            slide_width = self.prs.slide_width.inches
            box_width = (slide_width - 2) / 2 - 0.5
            
            for idx, item in enumerate(items[:2]):
                x_start = 0.5 if idx == 0 else slide_width/2 + 0.25
                
                # Box
                try:
                    box = slide.shapes.add_shape(
                        MSO_SHAPE.ROUNDED_RECTANGLE,
                        Inches(x_start), Inches(1.8),
                        Inches(box_width), Inches(4.8)
                    )
                    box.fill.solid()
                    box.fill.fore_color.rgb = RGBColor(30, 41, 59)
                    box.line.color.rgb = RGBColor(99, 102, 241)
                except:
                    pass
                
                # Title
                title_box = slide.shapes.add_textbox(
                    Inches(x_start + 0.2), Inches(2),
                    Inches(box_width - 0.4), Inches(0.6)
                )
                tf = title_box.text_frame
                p = tf.paragraphs[0]
                p.text = item.get('title', f'Option {idx + 1}')
                p.font.size = Pt(20)
                p.font.bold = True
                p.alignment = PP_ALIGN.CENTER
                
                # Pros
                y_pos = 2.8
                for pro in item.get('pros', [])[:3]:
                    pro_box = slide.shapes.add_textbox(
                        Inches(x_start + 0.2), Inches(y_pos),
                        Inches(box_width - 0.4), Inches(0.4)
                    )
                    tf = pro_box.text_frame
                    p = tf.paragraphs[0]
                    p.text = f"✓ {pro}"
                    p.font.size = Pt(12)
                    p.font.color.rgb = RGBColor(16, 185, 129)
                    y_pos += 0.45
                
                # Cons
                y_pos = max(y_pos, 4.5)
                for con in item.get('cons', [])[:3]:
                    con_box = slide.shapes.add_textbox(
                        Inches(x_start + 0.2), Inches(y_pos),
                        Inches(box_width - 0.4), Inches(0.4)
                    )
                    tf = con_box.text_frame
                    p = tf.paragraphs[0]
                    p.text = f"✗ {con}"
                    p.font.size = Pt(12)
                    p.font.color.rgb = RGBColor(239, 68, 68)
                    y_pos += 0.45
            
            # VS text
            vs_box = slide.shapes.add_textbox(
                Inches(slide_width/2 - 0.5), Inches(3.5),
                Inches(1), Inches(1)
            )
            tf = vs_box.text_frame
            p = tf.paragraphs[0]
            p.text = "VS"
            p.font.size = Pt(32)
            p.font.bold = True
            p.font.color.rgb = RGBColor(236, 72, 153)
            p.alignment = PP_ALIGN.CENTER
    
    def _fill_conclusion_slide(self, slide, slide_data: Dict):
        """Fill conclusion slide"""
        title = slide_data.get('title', 'Thank You')
        content_data = slide_data.get('content', {})
        bullets = content_data.get('bullet_points', [])
        cta = content_data.get('call_to_action', '')
        
        slide_width = self.prs.slide_width.inches
        
        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(2),
            Inches(slide_width - 1),
            Inches(1)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(44)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        
        # Bullets
        if bullets:
            y_pos = 3.5
            for bullet in bullets[:4]:
                bullet_box = slide.shapes.add_textbox(
                    Inches(2), Inches(y_pos),
                    Inches(slide_width - 4),
                    Inches(0.5)
                )
                tf = bullet_box.text_frame
                p = tf.paragraphs[0]
                p.text = f"• {bullet}"
                p.font.size = Pt(16)
                p.alignment = PP_ALIGN.CENTER
                y_pos += 0.5
        
        # CTA
        if cta:
            cta_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(6.3),
                Inches(slide_width - 1),
                Inches(0.5)
            )
            tf = cta_box.text_frame
            p = tf.paragraphs[0]
            p.text = cta
            p.font.size = Pt(16)
            p.font.color.rgb = RGBColor(236, 72, 153)
            p.alignment = PP_ALIGN.CENTER
    
    def _fill_image_slide(self, slide, slide_data: Dict):
        """Fill image placeholder slide"""
        title = slide_data.get('title', '')
        content_data = slide_data.get('content', {})
        image_desc = content_data.get('image_description', 'Image')
        
        slide_width = self.prs.slide_width.inches
        
        # Title
        if title:
            self._add_title(slide, title)
        
        # Image placeholder
        try:
            placeholder = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(1.5), Inches(1.8),
                Inches(slide_width - 3), Inches(4.5)
            )
            placeholder.fill.solid()
            placeholder.fill.fore_color.rgb = RGBColor(30, 41, 59)
            placeholder.line.color.rgb = RGBColor(100, 116, 139)
            placeholder.line.width = Pt(2)
        except:
            pass
        
        # Description
        desc_box = slide.shapes.add_textbox(
            Inches(1.5), Inches(3.5),
            Inches(slide_width - 3),
            Inches(1)
        )
        tf = desc_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"📷 {image_desc}"
        p.font.size = Pt(18)
        p.alignment = PP_ALIGN.CENTER
    
    def get_template_info(self) -> Dict:
        """Get template info"""
        if not self.prs:
            return {}
        
        return {
            'slide_width': self.prs.slide_width.inches,
            'slide_height': self.prs.slide_height.inches,
            'num_layouts': len(self.available_layouts),
            'layouts': [{'name': l['name'], 'index': l['index']} for l in self.available_layouts]
        }
