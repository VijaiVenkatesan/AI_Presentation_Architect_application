import streamlit as st
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
import json
import io
from groq import Groq
from typing import Dict, List, Any, Optional
import traceback

# ============================================================================
# CONFIGURATION & INITIALIZATION
# ============================================================================

st.set_page_config(
    page_title="AI Presentation Architect",
    page_icon="🎯",
    layout="wide"
)

# Initialize session state variables
if 'template_analyzed' not in st.session_state:
    st.session_state.template_analyzed = False
if 'layout_map' not in st.session_state:
    st.session_state.layout_map = {}
if 'ai_strategy' not in st.session_state:
    st.session_state.ai_strategy = None
if 'edited_strategy' not in st.session_state:
    st.session_state.edited_strategy = None
if 'available_models' not in st.session_state:
    st.session_state.available_models = []
if 'template_prs' not in st.session_state:
    st.session_state.template_prs = None

# ============================================================================
# GROQ API FUNCTIONS
# ============================================================================

def get_groq_client() -> Optional[Groq]:
    """Initialize and return Groq client using Streamlit secrets."""
    try:
        api_key = st.secrets.get("GROQ_API_KEY")
        if not api_key:
            st.error("⚠️ GROQ_API_KEY not found in Streamlit secrets.")
            return None
        return Groq(api_key=api_key)
    except Exception as e:
        st.error(f"❌ Error initializing Groq client: {str(e)}")
        return None

def fetch_available_models(client: Groq) -> List[str]:
    """Dynamically fetch available models from Groq API."""
    try:
        models = client.models.list()
        model_ids = [model.id for model in models.data if model.active]
        return sorted(model_ids)
    except Exception as e:
        st.error(f"❌ Error fetching models: {str(e)}")
        return []

def generate_presentation_strategy(
    client: Groq,
    model: str,
    topic: str,
    num_slides: int
) -> Optional[List[Dict[str, Any]]]:
    """Generate a complete presentation strategy using AI."""
    
    system_prompt = """You are an expert presentation strategist. Your task is to create a complete, 
professional presentation strategy in JSON format. You must generate diverse slide types including 
content slides, charts, and tables.

CRITICAL: You must respond with ONLY a valid JSON array. No explanations, no markdown, no code blocks.

Each slide object must have this exact structure:
{
    "type": "content" | "bar_chart" | "table",
    "title": "Slide Title Here",
    "data": {
        // For content type:
        "bullets": ["Point 1", "Point 2", "Point 3"]
        
        // For bar_chart type:
        "categories": ["Category 1", "Category 2", "Category 3"],
        "series": [{"name": "Series Name", "values": [10, 20, 30]}]
        
        // For table type:
        "headers": ["Column 1", "Column 2", "Column 3"],
        "rows": [["Data 1", "Data 2", "Data 3"], ["Data 4", "Data 5", "Data 6"]]
    }
}

Rules:
1. Create a mix of slide types (at least 1-2 charts and 1 table if num_slides > 5)
2. Keep bullet points concise (max 5 per slide)
3. Chart data must be realistic and relevant
4. Table data must be professional and formatted
5. Ensure titles are clear and descriptive
6. Output ONLY the JSON array, nothing else"""

    user_prompt = f"""Create a {num_slides}-slide presentation strategy about: {topic}

Generate a JSON array with exactly {num_slides} slide objects. Include a good mix of content slides, 
bar charts, and tables. Make it professional and comprehensive."""

    try:
        response = client.chat.completions.create(
            model=model,
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": user_prompt}
            ],
            temperature=0.7,
            max_tokens=4000
        )
        
        content = response.choices[0].message.content.strip()
        
        # Remove markdown code blocks if present
        if content.startswith("```"):
            content = content.split("```")[1]
            if content.startswith("json"):
                content = content[4:]
            content = content.strip()
        
        strategy = json.loads(content)
        
        # Validate structure
        if not isinstance(strategy, list):
            st.error("Invalid response: Expected a JSON array")
            return None
            
        for slide in strategy:
            if not all(k in slide for k in ['type', 'title', 'data']):
                st.error("Invalid slide structure in AI response")
                return None
        
        return strategy
        
    except json.JSONDecodeError as e:
        st.error(f"❌ Failed to parse AI response as JSON: {str(e)}")
        st.code(content, language="text")
        return None
    except Exception as e:
        st.error(f"❌ Error generating strategy: {str(e)}")
        return None

# ============================================================================
# TEMPLATE ANALYSIS FUNCTIONS
# ============================================================================

def analyze_template(prs: Presentation) -> Dict[str, Any]:
    """
    Deep analysis of PowerPoint template to identify and categorize layouts.
    
    Returns a dictionary with:
    - title_layout: The main title slide layout
    - content_layouts: List of layouts suitable for bullet content
    - chart_layouts: List of layouts with chart placeholders
    - table_layouts: List of layouts with table placeholders
    """
    
    layout_map = {
        'title_layout': None,
        'content_layouts': [],
        'chart_layouts': [],
        'table_layouts': [],
        'all_layouts': []
    }
    
    for idx, layout in enumerate(prs.slide_layouts):
        layout_info = {
            'index': idx,
            'name': layout.name,
            'has_title': False,
            'has_body': False,
            'has_chart': False,
            'has_table': False
        }
        
        # Analyze placeholders
        for shape in layout.placeholders:
            ph_type = shape.placeholder_format.type
            
            if ph_type == 1:  # Title placeholder
                layout_info['has_title'] = True
            elif ph_type == 2:  # Body/Content placeholder
                layout_info['has_body'] = True
            elif ph_type == 12:  # Chart placeholder
                layout_info['has_chart'] = True
            elif ph_type == 13:  # Table placeholder
                layout_info['has_table'] = True
        
        layout_map['all_layouts'].append(layout_info)
        
        # Categorize layouts
        # Title Slide: Has title, typically no body, usually first layout
        if idx == 0 or (layout_info['has_title'] and not layout_info['has_body'] 
                        and 'title' in layout.name.lower()):
            if layout_map['title_layout'] is None:
                layout_map['title_layout'] = idx
        
        # Chart Layout: Has chart placeholder
        if layout_info['has_chart']:
            layout_map['chart_layouts'].append(idx)
        
        # Table Layout: Has table placeholder
        if layout_info['has_table']:
            layout_map['table_layouts'].append(idx)
        
        # Content Layout: Has both title and body, but is NOT the title slide
        if (layout_info['has_title'] and layout_info['has_body'] 
            and idx != layout_map['title_layout']):
            layout_map['content_layouts'].append(idx)
    
    # Fallback: if no content layouts found, use any layout with title
    if not layout_map['content_layouts']:
        for info in layout_map['all_layouts']:
            if info['has_title'] and info['index'] != layout_map['title_layout']:
                layout_map['content_layouts'].append(info['index'])
    
    return layout_map

def display_layout_analysis(layout_map: Dict[str, Any]):
    """Display the template analysis results in the UI."""
    st.success("✅ Template analyzed successfully!")
    
    col1, col2, col3, col4 = st.columns(4)
    
    with col1:
        st.metric("Title Layout", 
                  f"#{layout_map['title_layout']}" if layout_map['title_layout'] is not None else "None")
    
    with col2:
        st.metric("Content Layouts", len(layout_map['content_layouts']))
    
    with col3:
        st.metric("Chart Layouts", len(layout_map['chart_layouts']))
    
    with col4:
        st.metric("Table Layouts", len(layout_map['table_layouts']))
    
    with st.expander("📋 Detailed Layout Information"):
        for info in layout_map['all_layouts']:
            st.write(f"**Layout {info['index']}: {info['name']}**")
            features = []
            if info['has_title']:
                features.append("Title")
            if info['has_body']:
                features.append("Body")
            if info['has_chart']:
                features.append("Chart")
            if info['has_table']:
                features.append("Table")
            st.write(f"  Features: {', '.join(features) if features else 'None detected'}")

# ============================================================================
# PRESENTATION BUILDING FUNCTIONS
# ============================================================================

def get_placeholder_by_type(slide, ph_type: int):
    """Get placeholder by type from a slide."""
    for shape in slide.placeholders:
        if shape.placeholder_format.type == ph_type:
            return shape
    return None

def add_title_slide(prs: Presentation, layout_idx: int, title: str, subtitle: str):
    """Add a title slide to the presentation."""
    slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
    
    # Try to set title
    title_shape = get_placeholder_by_type(slide, 1)  # Title type
    if title_shape and hasattr(title_shape, 'text_frame'):
        title_shape.text = title
    
    # Try to set subtitle
    subtitle_shape = get_placeholder_by_type(slide, 2)  # Body/Subtitle type
    if subtitle_shape and hasattr(subtitle_shape, 'text_frame'):
        subtitle_shape.text = subtitle

def add_content_slide(prs: Presentation, layout_idx: int, title: str, bullets: List[str]):
    """Add a content slide with bullet points."""
    slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
    
    # Set title
    title_shape = get_placeholder_by_type(slide, 1)
    if title_shape and hasattr(title_shape, 'text_frame'):
        title_shape.text = title
    
    # Set body content
    body_shape = get_placeholder_by_type(slide, 2)
    if body_shape and hasattr(body_shape, 'text_frame'):
        text_frame = body_shape.text_frame
        text_frame.clear()
        
        for i, bullet in enumerate(bullets):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            p.text = bullet
            p.level = 0

def add_chart_slide(prs: Presentation, layout_idx: int, title: str, chart_data: Dict[str, Any]):
    """Add a slide with a bar chart."""
    slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
    
    # Set title
    title_shape = get_placeholder_by_type(slide, 1)
    if title_shape and hasattr(title_shape, 'text_frame'):
        title_shape.text = title
    
    # Find chart placeholder
    chart_placeholder = get_placeholder_by_type(slide, 12)  # Chart type
    
    if chart_placeholder:
        # Prepare chart data
        chart_data_obj = CategoryChartData()
        chart_data_obj.categories = chart_data.get('categories', ['A', 'B', 'C'])
        
        for series in chart_data.get('series', []):
            chart_data_obj.add_series(
                series.get('name', 'Series'),
                series.get('values', [1, 2, 3])
            )
        
        # Add chart to placeholder
        chart = chart_placeholder.insert_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, chart_data_obj)
    else:
        # Fallback: add chart to slide if no placeholder
        x, y, cx, cy = Inches(1), Inches(2), Inches(8), Inches(4.5)
        
        chart_data_obj = CategoryChartData()
        chart_data_obj.categories = chart_data.get('categories', ['A', 'B', 'C'])
        
        for series in chart_data.get('series', []):
            chart_data_obj.add_series(
                series.get('name', 'Series'),
                series.get('values', [1, 2, 3])
            )
        
        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data_obj
        )

def add_table_slide(prs: Presentation, layout_idx: int, title: str, table_data: Dict[str, Any]):
    """Add a slide with a table."""
    slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
    
    # Set title
    title_shape = get_placeholder_by_type(slide, 1)
    if title_shape and hasattr(title_shape, 'text_frame'):
        title_shape.text = title
    
    headers = table_data.get('headers', [])
    rows = table_data.get('rows', [])
    
    # Find table placeholder
    table_placeholder = get_placeholder_by_type(slide, 13)  # Table type
    
    if table_placeholder:
        # Insert table into placeholder
        graphic_frame = table_placeholder.insert_table(
            rows=len(rows) + 1,
            cols=len(headers)
        )
        table = graphic_frame.table
    else:
        # Fallback: create table on slide
        x, y, cx, cy = Inches(1), Inches(2), Inches(8), Inches(4)
        shape = slide.shapes.add_table(
            len(rows) + 1, len(headers), x, y, cx, cy
        )
        table = shape.table
    
    # Populate headers
    for col_idx, header in enumerate(headers):
        cell = table.rows[0].cells[col_idx]
        cell.text = str(header)
        # Make header bold if possible
        if hasattr(cell.text_frame.paragraphs[0].runs[0], 'font'):
            cell.text_frame.paragraphs[0].runs[0].font.bold = True
    
    # Populate data rows
    for row_idx, row in enumerate(rows):
        for col_idx, cell_value in enumerate(row):
            if col_idx < len(table.rows[row_idx + 1].cells):
                table.rows[row_idx + 1].cells[col_idx].text = str(cell_value)

def build_presentation(
    template_prs: Presentation,
    layout_map: Dict[str, Any],
    main_title: str,
    topic: str,
    strategy: List[Dict[str, Any]]
) -> Optional[Presentation]:
    """Build the final presentation from the strategy."""
    
    try:
        # Create a new presentation from template
        prs = Presentation(io.BytesIO(template_prs._blob))
        
        # Add title slide
        if layout_map['title_layout'] is not None:
            add_title_slide(prs, layout_map['title_layout'], main_title, topic)
        
        # Track layout usage for cycling
        content_idx = 0
        chart_idx = 0
        table_idx = 0
        
        # Process each slide in the strategy
        for slide_spec in strategy:
            slide_type = slide_spec.get('type', 'content')
            title = slide_spec.get('title', 'Untitled')
            data = slide_spec.get('data', {})
            
            if slide_type == 'content':
                # Use content layouts with cycling
                if layout_map['content_layouts']:
                    layout_idx = layout_map['content_layouts'][content_idx % len(layout_map['content_layouts'])]
                    content_idx += 1
                    bullets = data.get('bullets', [])
                    add_content_slide(prs, layout_idx, title, bullets)
                else:
                    st.warning(f"⚠️ No content layouts available. Skipping slide: {title}")
            
            elif slide_type == 'bar_chart':
                # Use chart layouts with cycling
                if layout_map['chart_layouts']:
                    layout_idx = layout_map['chart_layouts'][chart_idx % len(layout_map['chart_layouts'])]
                    chart_idx += 1
                    add_chart_slide(prs, layout_idx, title, data)
                elif layout_map['content_layouts']:
                    # Fallback to content layout
                    layout_idx = layout_map['content_layouts'][content_idx % len(layout_map['content_layouts'])]
                    content_idx += 1
                    add_chart_slide(prs, layout_idx, title, data)
                else:
                    st.warning(f"⚠️ No suitable layouts for chart. Skipping slide: {title}")
            
            elif slide_type == 'table':
                # Use table layouts with cycling
                if layout_map['table_layouts']:
                    layout_idx = layout_map['table_layouts'][table_idx % len(layout_map['table_layouts'])]
                    table_idx += 1
                    add_table_slide(prs, layout_idx, title, data)
                elif layout_map['content_layouts']:
                    # Fallback to content layout
                    layout_idx = layout_map['content_layouts'][content_idx % len(layout_map['content_layouts'])]
                    content_idx += 1
                    add_table_slide(prs, layout_idx, title, data)
                else:
                    st.warning(f"⚠️ No suitable layouts for table. Skipping slide: {title}")
        
        return prs
        
    except Exception as e:
        st.error(f"❌ Error building presentation: {str(e)}")
        st.code(traceback.format_exc())
        return None

# ============================================================================
# MAIN APPLICATION UI
# ============================================================================

def main():
    st.title("🎯 True AI Presentation Architect")
    st.markdown("""
    **Transform your ideas into professional, branded presentations with AI.**
    
    Upload your template, define your topic, and let AI create a complete presentation 
    strategy with diverse content including charts and tables.
    """)
    
    # Initialize Groq client
    client = get_groq_client()
    if not client:
        st.stop()
    
    # Fetch available models if not already done
    if not st.session_state.available_models:
        with st.spinner("Fetching available AI models..."):
            st.session_state.available_models = fetch_available_models(client)
    
    # ========================================================================
    # STEP 1: UPLOAD TEMPLATE
    # ========================================================================
    
    st.header("📁 Step 1: Upload Your PowerPoint Template")
    
    uploaded_file = st.file_uploader(
        "Choose your .pptx template file",
        type=['pptx'],
        help="Upload your company's branded PowerPoint template"
    )
    
    if uploaded_file is not None:
        try:
            # Load presentation
            template_bytes = uploaded_file.read()
            prs = Presentation(io.BytesIO(template_bytes))
            
            # Store in session state (store the bytes, not the object)
            if st.session_state.template_prs is None:
                st.session_state.template_prs = Presentation(io.BytesIO(template_bytes))
                st.session_state.template_analyzed = False
            
            # Analyze template
            if not st.session_state.template_analyzed:
                with st.spinner("Analyzing template structure..."):
                    st.session_state.layout_map = analyze_template(prs)
                    st.session_state.template_analyzed = True
            
            display_layout_analysis(st.session_state.layout_map)
            
        except Exception as e:
            st.error(f"❌ Error loading template: {str(e)}")
            st.stop()
    else:
        st.info("👆 Please upload a PowerPoint template to begin")
        st.stop()
    
    # ========================================================================
    # STEP 2: DEFINE PRESENTATION DETAILS
    # ========================================================================
    
    st.header("✍️ Step 2: Define Your Presentation")
    
    col1, col2 = st.columns(2)
    
    with col1:
        main_title = st.text_input(
            "Main Presentation Title",
            placeholder="e.g., Q4 2024 Strategy Review",
            help="This will appear on the title slide"
        )
    
    with col2:
        topic = st.text_area(
            "Topic / Description",
            placeholder="e.g., Comprehensive review of Q4 performance, market analysis, and 2025 strategic initiatives",
            help="Describe what you want the presentation to cover",
            height=100
        )
    
    # ========================================================================
    # STEP 3: AI CONFIGURATION
    # ========================================================================
    
    st.header("🤖 Step 3: Configure AI Generation")
    
    col1, col2 = st.columns(2)
    
    with col1:
        if st.session_state.available_models:
            selected_model = st.selectbox(
                "Select AI Model",
                options=st.session_state.available_models,
                help="Choose the AI model to generate your presentation"
            )
        else:
            st.error("No models available. Please check your API configuration.")
            st.stop()
    
    with col2:
        num_slides = st.slider(
            "Number of Content Slides",
            min_value=3,
            max_value=20,
            value=8,
            help="How many slides should the AI generate? (excluding title slide)"
        )
    
    # ========================================================================
    # STEP 4: GENERATE STRATEGY
    # ========================================================================
    
    st.header("🎨 Step 4: Generate & Edit Presentation Strategy")
    
    if st.button("🚀 Generate AI Strategy", type="primary", use_container_width=True):
        if not main_title or not topic:
            st.error("⚠️ Please provide both a title and topic description")
        else:
            with st.spinner(f"AI is creating a {num_slides}-slide presentation strategy..."):
                strategy = generate_presentation_strategy(
                    client,
                    selected_model,
                    topic,
                    num_slides
                )
                
                if strategy:
                    st.session_state.ai_strategy = strategy
                    st.session_state.edited_strategy = json.loads(json.dumps(strategy))  # Deep copy
                    st.success(f"✅ Generated {len(strategy)} slides!")
                    st.rerun()
    
    # ========================================================================
    # DISPLAY & EDIT STRATEGY
    # ========================================================================
    
    if st.session_state.ai_strategy:
        st.subheader("📋 Review & Edit Your Presentation")
        
        st.info("💡 Expand each slide below to edit titles and content before building the final presentation.")
        
        for idx, slide in enumerate(st.session_state.edited_strategy):
            slide_type = slide.get('type', 'content')
            
            # Icon based on type
            if slide_type == 'bar_chart':
                icon = "📊"
            elif slide_type == 'table':
                icon = "📋"
            else:
                icon = "📄"
            
            with st.expander(f"{icon} Slide {idx + 1}: {slide.get('title', 'Untitled')} ({slide_type})"):
                # Edit title
                new_title = st.text_input(
                    "Title",
                    value=slide.get('title', ''),
                    key=f"title_{idx}"
                )
                st.session_state.edited_strategy[idx]['title'] = new_title
                
                # Edit content based on type
                if slide_type == 'content':
                    bullets = slide.get('data', {}).get('bullets', [])
                    st.write("**Bullet Points:**")
                    
                    new_bullets = []
                    for bullet_idx, bullet in enumerate(bullets):
                        new_bullet = st.text_input(
                            f"Bullet {bullet_idx + 1}",
                            value=bullet,
                            key=f"bullet_{idx}_{bullet_idx}"
                        )
                        new_bullets.append(new_bullet)
                    
                    # Option to add more bullets
                    if st.button(f"➕ Add Bullet Point", key=f"add_bullet_{idx}"):
                        new_bullets.append("New bullet point")
                    
                    st.session_state.edited_strategy[idx]['data']['bullets'] = new_bullets
                
                elif slide_type == 'bar_chart':
                    st.write("**Chart Data:**")
                    chart_data = slide.get('data', {})
                    
                    categories = chart_data.get('categories', [])
                    st.write(f"Categories: {', '.join(categories)}")
                    
                    for series_idx, series in enumerate(chart_data.get('series', [])):
                        st.write(f"**{series.get('name', 'Series')}:** {series.get('values', [])}")
                    
                    st.caption("📝 Chart data editing in UI coming soon. Edit JSON below for now.")
                
                elif slide_type == 'table':
                    st.write("**Table Data:**")
                    table_data = slide.get('data', {})
                    
                    headers = table_data.get('headers', [])
                    rows = table_data.get('rows', [])
                    
                    st.write(f"**Headers:** {', '.join(headers)}")
                    st.write(f"**Rows:** {len(rows)}")
                    
                    for row_idx, row in enumerate(rows[:3]):  # Show first 3 rows
                        st.write(f"Row {row_idx + 1}: {', '.join(map(str, row))}")
                    
                    if len(rows) > 3:
                        st.caption(f"... and {len(rows) - 3} more rows")
                    
                    st.caption("📝 Table data editing in UI coming soon. Edit JSON below for now.")
        
        # Advanced: Raw JSON editing
        with st.expander("🔧 Advanced: Edit Raw JSON"):
            edited_json = st.text_area(
                "Edit the complete strategy as JSON",
                value=json.dumps(st.session_state.edited_strategy, indent=2),
                height=400
            )
            
            if st.button("Apply JSON Changes"):
                try:
                    new_strategy = json.loads(edited_json)
                    st.session_state.edited_strategy = new_strategy
                    st.success("✅ JSON changes applied!")
                    st.rerun()
                except json.JSONDecodeError as e:
                    st.error(f"❌ Invalid JSON: {str(e)}")
        
        # ====================================================================
        # STEP 5: BUILD FINAL PRESENTATION
        # ====================================================================
        
        st.header("🎯 Step 5: Build Your Presentation")
        
        if st.button("🏗️ Build Final Presentation", type="primary", use_container_width=True):
            with st.spinner("Building your presentation..."):
                final_prs = build_presentation(
                    st.session_state.template_prs,
                    st.session_state.layout_map,
                    main_title,
                    topic,
                    st.session_state.edited_strategy
                )
                
                if final_prs:
                    # Save to bytes
                    prs_bytes = io.BytesIO()
                    final_prs.save(prs_bytes)
                    prs_bytes.seek(0)
                    
                    st.success("✅ Presentation built successfully!")
                    
                    # Download button
                    st.download_button(
                        label="📥 Download Presentation",
                        data=prs_bytes,
                        file_name=f"{main_title.replace(' ', '_')}.pptx",
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                        use_container_width=True
                    )
                    
                    # Summary
                    st.info(f"""
                    **Presentation Summary:**
                    - Title: {main_title}
                    - Total Slides: {len(st.session_state.edited_strategy) + 1} (including title slide)
                    - Content Slides: {sum(1 for s in st.session_state.edited_strategy if s['type'] == 'content')}
                    - Charts: {sum(1 for s in st.session_state.edited_strategy if s['type'] == 'bar_chart')}
                    - Tables: {sum(1 for s in st.session_state.edited_strategy if s['type'] == 'table')}
                    """)

if __name__ == "__main__":
    main()
