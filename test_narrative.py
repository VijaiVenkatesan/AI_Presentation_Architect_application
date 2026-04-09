
import os
import sys
from pptx import Presentation

# Add path
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))

from core.ppt_generator import PptGenerator
from core.template_parser import PptxTemplateParser

def test_narrative_branding():
    template_path = r"D:\Antigravity_Projects\AI_Presentation_Architecture_App\assets\Datamatics-Presentation-Template-Light-Test.pptx"
    out_path = r"C:\Users\vijai.v\Downloads\test_narrative_branding.pptx"
    
    slides = [
        {"title": "Title slide", "bullet_points": []},
        {"title": "Narrative Slide 1", "subtitle": "Should use Layout 2", "bullet_points": ["Point A", "Point B"]},
        {"title": "Narrative Slide 2", "subtitle": "Should use Layout 2", "bullet_points": ["Point C", "Point D"]},
        {"title": "Thank You", "bullet_points": []}
    ]
    
    with open(template_path, "rb") as f:
        template_bytes = f.read()
    
    parser = PptxTemplateParser(template_bytes)
    profile = parser.parse()
    
    generator = PptGenerator(template_bytes=template_bytes, template_profile=profile)
    pptx_bytes = generator.generate(slides)
    
    with open(out_path, "wb") as f:
        f.write(pptx_bytes)
    
    prs = Presentation(out_path)
    print(f"Slide count: {len(prs.slides)}")
    for i, slide in enumerate(prs.slides):
        print(f"Slide {i+1} Layout: {slide.slide_layout.name}")

if __name__ == "__main__":
    test_narrative_branding()
