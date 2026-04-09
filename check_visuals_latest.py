
from pptx import Presentation
import os

def check(path):
    if not os.path.exists(path):
        print(f"File not found: {path}")
        return
    print(f"\nChecking: {os.path.basename(path)}")
    prs = Presentation(path)
    for i, slide in enumerate(prs.slides):
        stypes = [str(s.shape_type) for s in slide.shapes]
        has_v = 'CHART (3)' in stypes or 'TABLE (19)' in stypes
        print(f"Slide {i+1}: {stypes} {'(VISUAL!)' if has_v else ''}")

check(r'C:\Users\vijai.v\Downloads\presentation_20260330_132232.pptx')
