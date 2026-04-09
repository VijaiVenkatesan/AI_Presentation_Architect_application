
from pptx import Presentation
from pptx.util import Inches, Emu
import os

def analyze_pptx(path):
    if not os.path.exists(path):
        print(f"Error: File not found at {path}")
        return

    # Check if we can open it
    try:
        prs = Presentation(path)
    except Exception as e:
        print(f"Error opening PPTX: {e}")
        return

    print(f"Slide Width: {prs.slide_width.inches:.2f}\"")
    print(f"Slide Height: {prs.slide_height.inches:.2f}\"")
    print(f"Slide count: {len(prs.slides)}")
    
    for i, slide in enumerate(prs.slides):
        print(f"\nSlide {i+1}:")
        print(f"  Layout index: {prs.slide_layouts.index(slide.slide_layout)}")
        print(f"  Layout name: {slide.slide_layout.name}")
        
        shapes = slide.shapes
        print(f"  Shape count: {len(shapes)}")
        
        for j, shape in enumerate(shapes):
            try:
                left = shape.left.inches
                top = shape.top.inches
                width = shape.width.inches
                height = shape.height.inches
                
                print(f"  Shape {j+1}: {shape.name}")
                print(f"    Position: Left={left:.2f}\", Top={top:.2f}\"")
                print(f"    Size: Width={width:.2f}\", Height={height:.2f}\"")
                
                if hasattr(shape, "text") and shape.text.strip():
                    print(f"    Text: \"{shape.text[:50]}...\"")
            except Exception as e:
                print(f"    Shape {j+1}: {shape.name} (Error reading properties: {e})")
        
        if i >= 2: # Just first 3 slides
           break

if __name__ == "__main__":
    path = r"C:\Users\vijai.v\Downloads\presentation_20260330_114628.pptx"
    analyze_pptx(path)
