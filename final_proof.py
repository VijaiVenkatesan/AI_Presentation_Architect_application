
import os
from pptx import Presentation
from pptx.util import Inches

def perform_evidence_check():
    FILE_PATH = "true_gamma_test.pptx"
    if not os.path.exists(FILE_PATH):
        print(f"ERROR: {FILE_PATH} not found. Run the generator first.")
        return

    prs = Presentation(FILE_PATH)
    print("--- FINAL EVIDENCE CHECK (TRUE-GAMMA v4.0) ---")
    print(f"Total Slides: {len(prs.slides)}")
    
    # 1. Sequence Check
    for i, slide in enumerate(prs.slides):
        print(f"\n[Slide {i}] Layout: {slide.slide_layout.name}")
        
        # 2. Shape Count (Proof of Sterile Injection)
        print(f"  Shape Count: {len(slide.shapes)}")
        for j, shape in enumerate(slide.shapes):
            top_in = shape.top / Inches(1)
            text = ""
            if hasattr(shape, 'text_frame'):
                 if shape.text_frame and shape.text_frame.paragraphs:
                     text = shape.text_frame.text[:30]
            print(f"    - {shape.name} | Top: {top_in:.2f}in | Text: {text}")

    # 3. Nuclear Proof: Check if 'Picture' or 'Group' from Layout survived
    unexpected_pests = []
    for i, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if any(x in shape.name for x in ["Picture", "Group", "Freeform", "Oval", "Rectangle"]):
                 # Exclude text boxes or added markers
                 is_content = False
                 if hasattr(shape, 'text_frame'):
                      if "AI VISUAL" in shape.text_frame.text: is_content = True
                      if len(shape.text_frame.text.strip()) > 0: is_content = True
                 
                 if not is_content:
                      unexpected_pests.append((i, shape.name))

    print("\n--- NUCLEAR VERIFICATION ---")
    if not unexpected_pests:
        print("PASS: Zero 'Ghost' Laptops, Maps, or Groups detected.")
    else:
        print(f"FAIL: Detected {len(unexpected_pests)} stray elements.")
        for slide_idx, name in unexpected_pests:
            print(f"    - Slide {slide_idx}: {name}")

    print("\n--- BRANDING FORTRESS CHECK ---")
    high_shapes = []
    for i, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.top / Inches(1) < 1.4 and i < 2: 
                 high_shapes.append((i, shape.name, shape.top / Inches(1)))
    
    if not high_shapes:
        print("PASS: Branding Zone is 100% CLEAR for Red Bar/Logo.")
    else:
        print(f"WARNING: {len(high_shapes)} shapes in branding zone.")

if __name__ == "__main__":
    perform_evidence_check()
