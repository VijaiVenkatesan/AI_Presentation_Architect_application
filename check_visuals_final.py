
from pptx import Presentation
import os

def check(path):
    print(f"\nChecking: {os.path.basename(path)}")
    prs = Presentation(path)
    for i, slide in enumerate(prs.slides):
        stypes = [str(s.shape_type) for s in slide.shapes]
        if 'CHART (3)' in stypes or 'TABLE (19)' in stypes:
            print(f"Slide {i+1}: {stypes}")

check(r'C:\Users\vijai.v\Downloads\presentation_20260330_125520.pptx')
check(r'C:\Users\vijai.v\Downloads\presentation_20260330_125650.pptx')
