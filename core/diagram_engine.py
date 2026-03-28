from pptx.util import Inches

def add_diagram(slide, diagram_type):
    if diagram_type == "flow":
        slide.shapes.add_textbox(Inches(1), Inches(2), Inches(2), Inches(1)).text = "Start"
        slide.shapes.add_textbox(Inches(4), Inches(2), Inches(2), Inches(1)).text = "Process"
        slide.shapes.add_textbox(Inches(7), Inches(2), Inches(2), Inches(1)).text = "End"
