from pptx import Presentation
import copy

def clone_slide(prs, index):
    source = prs.slides[index]
    blank = prs.slide_layouts[6]
    new_slide = prs.slides.add_slide(blank)

    for shape in source.shapes:
        new_el = copy.deepcopy(shape.element)
        new_slide.shapes._spTree.insert_element_before(new_el, 'p:extLst')

    return new_slide


def replace_text(slide, data):
    for shape in slide.shapes:
        if shape.has_text_frame:
            shape.text = data.get("title", "") + "\n" + "\n".join(data.get("bullet_points", []))
