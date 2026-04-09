"""
Enterprise AI Presentation Architect — Template Parser
Parses PPTX files and images to extract layout, styling, and branding metadata.
"""

import io
import logging
from copy import deepcopy
from typing import Dict, List, Optional, Tuple

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import PP_PLACEHOLDER

logger = logging.getLogger("PresentationArchitect")


# ─── PPTX Template Parser ───────────────────────────────────────────────────────

class PptxTemplateParser:
    """
    Parses uploaded PPTX templates to extract complete styling metadata.
    Enables faithful reproduction of template look and feel.
    """

    def __init__(self, pptx_bytes: bytes):
        self.pptx_bytes = pptx_bytes
        self.presentation = None
        self.profile = {}

    def parse(self) -> Dict:
        """
        Full template parsing pipeline.
        Returns a comprehensive TemplateProfile dictionary.
        """
        try:
            self.presentation = Presentation(io.BytesIO(self.pptx_bytes))
            logger.info("PPTX template loaded successfully")
        except Exception as e:
            logger.error(f"Failed to load PPTX: {e}")
            return self._empty_profile("Failed to parse PPTX template")

        self.profile = {
            "source_type": "pptx",
            "slide_width": self.presentation.slide_width,
            "slide_height": self.presentation.slide_height,
            "layouts": [],
            "color_scheme": {},
            "fonts": {},
            "logo_info": None,
            "background_color": None,
            "placeholder_map": {},
            "total_slides": len(self.presentation.slides),
            "theme_name": "",
            "slide_details": [],
        }

        self._extract_layouts()
        self._extract_color_scheme()
        self._extract_fonts()
        self._extract_slide_details()
        self._detect_logos()

        logger.info(
            f"Template parsed: {len(self.profile['layouts'])} layouts, "
            f"{self.profile['total_slides']} slides"
        )
        return self.profile

    def _extract_layouts(self):
        """Extract all slide layouts from slide masters."""
        try:
            for master_idx, master in enumerate(self.presentation.slide_masters):
                for layout_idx, layout in enumerate(master.slide_layouts):
                    layout_info = {
                        "master_index": master_idx,
                        "layout_index": layout_idx,
                        "name": layout.name or f"Layout {layout_idx}",
                        "placeholders": [],
                        "has_title": False,
                        "has_body": False,
                        "has_subtitle": False,
                    }

                    for ph in layout.placeholders:
                        ph_info = {
                            "idx": ph.placeholder_format.idx,
                            "type": str(ph.placeholder_format.type),
                            "name": ph.name,
                            "left": ph.left,
                            "top": ph.top,
                            "width": ph.width,
                            "height": ph.height,
                        }

                        # Classify placeholder using official Enums (Type is more reliable than IDX)
                        ph_type = ph.placeholder_format.type
                        if ph_type in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE):
                            layout_info["has_title"] = True
                            ph_info["role"] = "title"
                        elif ph_type == PP_PLACEHOLDER.BODY:
                            layout_info["has_body"] = True
                            ph_info["role"] = "body"
                        elif ph_type in (PP_PLACEHOLDER.SUBTITLE, 4):
                            layout_info["has_subtitle"] = True
                            ph_info["role"] = "subtitle"
                        elif ph_type == PP_PLACEHOLDER.OBJECT:
                            # Object is often the main content box in custom templates
                            ph_info["role"] = "body"
                        else:
                            ph_info["role"] = "other"

                        # Extract font info from placeholder
                        try:
                            if ph.text_frame:
                                for para in ph.text_frame.paragraphs:
                                    if para.font:
                                        ph_info["font_name"] = para.font.name
                                        ph_info["font_size"] = (
                                            para.font.size.pt if para.font.size else None
                                        )
                                        ph_info["font_bold"] = para.font.bold
                                        ph_info["font_italic"] = para.font.italic
                                        if para.font.color and para.font.color.rgb:
                                            ph_info["font_color"] = str(para.font.color.rgb)
                                        break
                        except Exception:
                            pass

                        layout_info["placeholders"].append(ph_info)

                    self.profile["layouts"].append(layout_info)
                    self.profile["placeholder_map"][layout_idx] = layout_info["placeholders"]

        except Exception as e:
            logger.warning(f"Error extracting layouts: {e}")

    def _extract_color_scheme(self):
        """Extract theme color scheme."""
        try:
            theme = self.presentation.slide_masters[0].element
            # Try to get theme colors from XML
            ns = {
                'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
            }

            # Search for color scheme in theme XML
            theme_elements = theme.findall('.//a:clrScheme', ns)
            if not theme_elements:
                # Try from the theme part directly
                try:
                    theme_part = self.presentation.slide_masters[0].slide_layouts[0].part
                    theme_name = "Default"
                except Exception:
                    theme_name = "Default"
                self.profile["theme_name"] = theme_name

            # Extract colors from first slide if available
            if self.presentation.slides:
                slide = self.presentation.slides[0]
                colors_found = set()
                for shape in slide.shapes:
                    try:
                        if shape.has_text_frame:
                            for para in shape.text_frame.paragraphs:
                                for run in para.runs:
                                    if run.font.color and run.font.color.rgb:
                                        colors_found.add(str(run.font.color.rgb))
                    except Exception:
                        pass

                if colors_found:
                    self.profile["color_scheme"]["text_colors"] = list(colors_found)

        except Exception as e:
            logger.warning(f"Error extracting colors: {e}")

    def _extract_fonts(self):
        """Extract font information from slides."""
        try:
            fonts_used = {}
            title_fonts = []
            body_fonts = []

            for slide in self.presentation.slides:
                for shape in slide.shapes:
                    if not shape.has_text_frame:
                        continue
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            font_name = run.font.name
                            font_size = run.font.size.pt if run.font.size else None
                            if font_name:
                                if font_name not in fonts_used:
                                    fonts_used[font_name] = {
                                        "name": font_name,
                                        "sizes": [],
                                        "count": 0
                                    }
                                fonts_used[font_name]["count"] += 1
                                if font_size:
                                    fonts_used[font_name]["sizes"].append(font_size)

                        # Detect title vs body fonts
                        if para.level == 0 and para.runs:
                            fn = para.runs[0].font.name
                            if fn:
                                title_fonts.append(fn)
                        elif para.runs:
                            fn = para.runs[0].font.name
                            if fn:
                                body_fonts.append(fn)

            self.profile["fonts"] = {
                "all_fonts": fonts_used,
                "primary_title_font": max(set(title_fonts), key=title_fonts.count) if title_fonts else None,
                "primary_body_font": max(set(body_fonts), key=body_fonts.count) if body_fonts else None,
            }

        except Exception as e:
            logger.warning(f"Error extracting fonts: {e}")

    def _extract_slide_details(self):
        """Extract details from each existing slide."""
        try:
            for idx, slide in enumerate(self.presentation.slides):
                slide_info = {
                    "index": idx,
                    "layout_name": slide.slide_layout.name if slide.slide_layout else "Unknown",
                    "layout_index": self._find_layout_index(slide.slide_layout),
                    "shapes_count": len(slide.shapes),
                    "has_images": False,
                    "has_charts": False,
                    "has_tables": False,
                    "placeholder_texts": {},
                }

                for shape in slide.shapes:
                    if shape.shape_type == 13:  # Picture
                        slide_info["has_images"] = True
                    if shape.has_chart:
                        slide_info["has_charts"] = True
                    if shape.has_table:
                        slide_info["has_tables"] = True
                    if shape.has_text_frame and hasattr(shape, "placeholder_format"):
                        try:
                            if shape.placeholder_format:
                                idx_val = shape.placeholder_format.idx
                                slide_info["placeholder_texts"][idx_val] = shape.text
                        except Exception:
                            pass

                self.profile["slide_details"].append(slide_info)

        except Exception as e:
            logger.warning(f"Error extracting slide details: {e}")

    def _detect_logos(self):
        """Detect logo images in the template."""
        try:
            for slide in self.presentation.slides:
                for shape in slide.shapes:
                    if shape.shape_type == 13:  # Picture
                        # Heuristic: headers/logos are in top corner or very small
                        w, h, l, t = shape.width, shape.height, shape.left, shape.top
                        sw, sh = self.profile["slide_width"], self.profile["slide_height"]

                        is_header = t < sh * 0.15
                        is_logo = (w < sw * 0.2 and h < sh * 0.2)
                        
                        if is_header and is_logo:
                             # Track the 'Brand Floor' — No AI content can start above this line
                             current_floor = self.profile.get("brand_safe_floor", Emu(Pt(50)))
                             self.profile["brand_safe_floor"] = max(current_floor, t + h + Emu(Pt(20)))

                             if not self.profile["logo_info"]:
                                 self.profile["logo_info"] = {"left":l, "top":t, "width":w, "height":h, "name":shape.name}
                                 logger.info(f"Detected brand logo/header at top: {shape.name}")

        except Exception as e:
            logger.warning(f"Error detecting logos: {e}")

    def _find_layout_index(self, slide_layout) -> int:
        """Find the index of a slide layout within all layouts."""
        try:
            for master in self.presentation.slide_masters:
                for idx, layout in enumerate(master.slide_layouts):
                    if layout.name == slide_layout.name:
                        return idx
        except Exception:
            pass
        return 0

    def _empty_profile(self, error: str = "") -> Dict:
        """Return an empty profile with error message."""
        return {
            "source_type": "pptx",
            "error": error,
            "layouts": [],
            "slide_width": 12192000,
            "slide_height": 6858000,
            "color_scheme": {},
            "fonts": {},
            "logo_info": None,
            "background_color": None,
            "placeholder_map": {},
            "total_slides": 0,
            "theme_name": "",
            "slide_details": [],
        }

    def get_best_layout_for_content(self, content: Dict) -> int:
        """
        Determine the best layout index for given slide content.
        Uses heuristics based on content type.
        """
        layouts = self.profile.get("layouts", [])
        if not layouts:
            return 0

        has_chart = content.get("chart_data") is not None
        has_table = content.get("table_data") is not None
        has_bullets = bool(content.get("bullet_points"))
        slide_num = content.get("slide_number", 1)

        # Title slide
        if slide_num == 1:
            for i, layout in enumerate(layouts):
                name_lower = layout.get("name", "").lower()
                if "title" in name_lower and ("slide" in name_lower or "only" not in name_lower):
                    return i
            return 0

        # Content with chart/table — look for "blank" or "content" layout
        if has_chart or has_table:
            for i, layout in enumerate(layouts):
                name_lower = layout.get("name", "").lower()
                if "blank" in name_lower:
                    return i
            for i, layout in enumerate(layouts):
                name_lower = layout.get("name", "").lower()
                if "content" in name_lower:
                    return i

        # Bullet content — look for body placeholder layouts
        if has_bullets:
            for i, layout in enumerate(layouts):
                if layout.get("has_title") and layout.get("has_body"):
                    return i

        # Default: use layout index 1 (usually "Title and Content")
        return min(1, len(layouts) - 1)


# ─── Image Template Parser ──────────────────────────────────────────────────────

class ImageTemplateParser:
    """
    Parses uploaded slide images to infer layout structure.
    Uses OCR for text extraction and basic color analysis.
    """

    def __init__(self, image_bytes: bytes, filename: str = ""):
        self.image_bytes = image_bytes
        self.filename = filename

    def parse(self) -> Dict:
        """
        Parse an image to extract layout information.
        Returns profile dictionary with inferred layout data.
        """
        profile = {
            "source_type": "image",
            "slide_width": 12192000,
            "slide_height": 6858000,
            "layouts": [],
            "color_scheme": {},
            "fonts": {"primary_title_font": "Calibri", "primary_body_font": "Calibri"},
            "logo_info": None,
            "background_color": None,
            "placeholder_map": {},
            "total_slides": 1,
            "theme_name": "Image-based",
            "slide_details": [],
            "ocr_text": "",
            "inferred_regions": [],
        }

        try:
            from PIL import Image
            img = Image.open(io.BytesIO(self.image_bytes))
            width, height = img.size

            # Calculate aspect ratio to determine slide dimensions
            aspect = width / height
            if abs(aspect - 16/9) < abs(aspect - 4/3):
                # 16:9 widescreen
                profile["slide_width"] = 12192000
                profile["slide_height"] = 6858000
            else:
                # 4:3 standard
                profile["slide_width"] = 9144000
                profile["slide_height"] = 6858000

            # Extract dominant colors
            profile["color_scheme"] = self._extract_colors(img)
            profile["background_color"] = profile["color_scheme"].get("background")

            # OCR text extraction
            profile["ocr_text"] = self._extract_text(img)

            # Infer layout regions
            profile["inferred_regions"] = self._infer_regions(img, profile["ocr_text"])

            # Create a default layout entry
            profile["layouts"] = [{
                "name": "Inferred Layout",
                "layout_index": 0,
                "has_title": True,
                "has_body": True,
                "placeholders": [
                    {"idx": 0, "role": "title",
                     "left": int(profile["slide_width"] * 0.05),
                     "top": int(profile["slide_height"] * 0.05),
                     "width": int(profile["slide_width"] * 0.9),
                     "height": int(profile["slide_height"] * 0.15)},
                    {"idx": 1, "role": "body",
                     "left": int(profile["slide_width"] * 0.05),
                     "top": int(profile["slide_height"] * 0.25),
                     "width": int(profile["slide_width"] * 0.9),
                     "height": int(profile["slide_height"] * 0.65)},
                ]
            }]

        except ImportError as e:
            logger.warning(f"Image processing library not available: {e}")
        except Exception as e:
            logger.warning(f"Image parsing error: {e}")

        return profile

    def _extract_colors(self, img) -> Dict:
        """Extract dominant colors from image."""
        try:
            # Resize for speed
            small = img.resize((100, 100))
            pixels = list(small.getdata())

            if not pixels:
                return {}

            # Get background (most common color from corners)
            corners = []
            w, h = small.size
            for x, y in [(0, 0), (w-1, 0), (0, h-1), (w-1, h-1)]:
                corners.append(small.getpixel((x, y)))

            # Convert to RGB if needed
            if isinstance(corners[0], int):
                bg_color = f"#{corners[0]:02x}{corners[0]:02x}{corners[0]:02x}"
            elif len(corners[0]) >= 3:
                r = sum(c[0] for c in corners) // 4
                g = sum(c[1] for c in corners) // 4
                b = sum(c[2] for c in corners) // 4
                bg_color = f"#{r:02x}{g:02x}{b:02x}"
            else:
                bg_color = "#ffffff"

            # Simple color frequency analysis
            from collections import Counter
            # Quantize colors
            quantized = []
            for p in pixels:
                if isinstance(p, int):
                    quantized.append((p // 32 * 32,) * 3)
                elif len(p) >= 3:
                    quantized.append((p[0] // 32 * 32, p[1] // 32 * 32, p[2] // 32 * 32))

            color_counts = Counter(quantized).most_common(5)
            accent_colors = [f"#{c[0]:02x}{c[1]:02x}{c[2]:02x}" for (c, _) in color_counts if c]

            return {
                "background": bg_color,
                "accent_colors": accent_colors,
                "text_colors": ["#000000", "#333333"]
            }

        except Exception as e:
            logger.warning(f"Color extraction error: {e}")
            return {"background": "#ffffff", "accent_colors": [], "text_colors": ["#000000"]}

    def _extract_text(self, img) -> str:
        """Extract text from image using OCR."""
        try:
            import pytesseract
            text = pytesseract.image_to_string(img)
            return text.strip()
        except ImportError:
            logger.info("pytesseract not available — OCR text extraction skipped")
            return ""
        except Exception as e:
            logger.warning(f"OCR extraction error: {e}")
            return ""

    def _infer_regions(self, img, ocr_text: str) -> List[Dict]:
        """Infer layout regions from image analysis."""
        regions = []
        w, h = img.size

        # Basic region inference based on standard slide layout
        regions.append({
            "type": "title",
            "x": int(w * 0.05),
            "y": int(h * 0.05),
            "width": int(w * 0.9),
            "height": int(h * 0.15),
        })

        regions.append({
            "type": "body",
            "x": int(w * 0.05),
            "y": int(h * 0.25),
            "width": int(w * 0.9),
            "height": int(h * 0.65),
        })

        regions.append({
            "type": "footer",
            "x": int(w * 0.05),
            "y": int(h * 0.92),
            "width": int(w * 0.9),
            "height": int(h * 0.06),
        })

        return regions
