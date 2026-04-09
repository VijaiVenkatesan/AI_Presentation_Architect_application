"""
Enterprise AI Presentation Architect — PPT Generator (True-Blanking v5.0)
The ABSOLUTE LOCKDOWN: Forces Layout 4 (Blank) for every slide to 
eliminate template 'ghost' shapes (Laptops/Maps) while keeping Master Branding.
"""

import io
import logging
from typing import Dict, List, Optional

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.chart import XL_CHART_TYPE
from pptx.oxml.ns import qn

logger = logging.getLogger("PresentationArchitect")

DEFAULT_THEME = {
    "title_font": "Calibri", "body_font": "Calibri",
    "title_size": Pt(34), "subtitle_size": Pt(20),
    "body_size": Pt(20), "bullet_size": Pt(18),
    "title_color": RGBColor(0x1A, 0x1A, 0x2E),
    "body_color": RGBColor(0x33, 0x33, 0x33),
}

class PptGenerator:
    def __init__(self, template_bytes=None, template_profile=None):
        self.template_bytes = template_bytes
        self.template_profile = template_profile or {}
        self.prs = None
        self._theme = dict(DEFAULT_THEME)
        
        # Style Inference
        fonts = self.template_profile.get("fonts", {})
        if fonts.get("primary_title_font"): self._theme["title_font"] = fonts["primary_title_font"]
        if fonts.get("primary_body_font"): self._theme["body_font"] = fonts["primary_body_font"]
        
        colors = self.template_profile.get("colors", {})
        if colors.get("primary"):
             try:
                 c = colors["primary"].lstrip('#')
                 self._theme["title_color"] = RGBColor(int(c[0:2], 16), int(c[2:4], 16), int(c[4:6], 16))
             except: pass

    def generate(self, slides_content: List[Dict], progress_callback=None) -> bytes:
        if self.template_bytes:
            self.prs = Presentation(io.BytesIO(self.template_bytes))
            self._remove_existing_slides()
        else:
            self.prs = Presentation()
            self.prs.slide_width, self.prs.slide_height = Emu(12192000), Emu(6858000)

        num_slides = len(slides_content)
        for i, content in enumerate(slides_content):
            if progress_callback: progress_callback((i+1)/num_slides, f"🔒 Smart Layout Generation: Slide {i+1}/{num_slides}...")
            
            if i == 0:
                # Title slide (Layout 0)
                self._add_template_slide(content, "title")
            elif i == num_slides - 1:
                # Thank you slide (Layout 20 or similar)
                self._add_template_slide(content, "thank_you")
            else:
                # Content slide (Standard Layout + Sterile Stack)
                self._add_content_slide(content)

        output = io.BytesIO()
        self.prs.save(output)
        output.seek(0)
        return output.getvalue()

    def _remove_existing_slides(self):
        xml_slides = self.prs.slides._sldIdLst
        for slide_id in list(xml_slides):
            rId = slide_id.get(qn('r:id'))
            self.prs.part.drop_rel(rId)
            xml_slides.remove(slide_id)

    def _add_template_slide(self, content: Dict, slide_type: str):
        """Adds a slide using its template layout and fills native placeholders."""
        layout_idx = 0
        if slide_type == "title":
            layout_idx = 0
        elif slide_type == "thank_you":
            # Search for 'Thank You' in layout names, default to 20 or 0
            layout_idx = 0
            for i, l in enumerate(self.prs.slide_layouts):
                if "thank" in l.name.lower(): layout_idx = i; break
            if layout_idx == 0 and len(self.prs.slide_layouts) > 20: layout_idx = 20

        slide = self.prs.slides.add_slide(self.prs.slide_layouts[layout_idx])
        self._fill_placeholders(slide, content)

    def _fill_placeholders(self, slide, content):
        """Fills title and subtitle placeholders in a template layout."""
        for ph in slide.placeholders:
            ph_type = ph.placeholder_format.type
            if ph_type in (1, 3): # TITLE or CENTER_TITLE
                ph.text = content.get("title", "")
            elif ph_type in (4, 10): # SUBTITLE or BODY
                ph.text = content.get("subtitle", "") or "\n".join(content.get("bullet_points", []))

    def _add_content_slide(self, content: Dict):
        """Adds a content slide. Bifurcates between Nuclear-Blank (Visuals) and Template-Native (Narrative)."""
        has_visuals = bool(content.get("chart_data") or content.get("table_data") or content.get("image_prompt"))
        
        layout = None
        is_nuclear = False

        if has_visuals:
            # RESTORE HAPPY STATE: Use Layout 4 (Blank) for heavy visual slides
            try:
                layout = self.prs.slide_layouts[4]
                is_nuclear = True
            except:
                layout = self.prs.slide_layouts[0]
        else:
            # TEMPLATE ADHERENCE: Use Layout 2 (Title and Content) for narrative decks
            try:
                if len(self.prs.slide_layouts) > 2:
                    layout = self.prs.slide_layouts[2]
                else:
                    layout = self.prs.slide_layouts[0]
            except:
                layout = self.prs.slide_layouts[0]

        # Search by name fallback
        if not layout or ("title" in layout.name.lower() and "content" not in layout.name.lower()):
            for l in self.prs.slide_layouts:
                if "content" in l.name.lower(): layout = l; break
            if not layout:
                for l in self.prs.slide_layouts:
                    if "blank" in l.name.lower(): layout = l; is_nuclear = True; break

        slide = self.prs.slides.add_slide(layout or self.prs.slide_layouts[0])
        
        # SAFE CLEARING
        for shape in list(slide.shapes):
            if hasattr(shape, "placeholder_format"):
                try: shape.element.getparent().remove(shape.element)
                except: pass
            elif is_nuclear:
                # Nuclear clear for Visual slides (User's Happy State)
                try: shape.element.getparent().remove(shape.element)
                except: pass
            else:
                # Gentle clear for Narrative (preserves footers/logo)
                if shape.top < Emu(Inches(7.0)) and shape.top > Emu(Inches(2.0)):
                    try: shape.element.getparent().remove(shape.element)
                    except: pass

        # Pass is_nuclear flag to avoid loop shredding
        self._build_sterile_stack(slide, content, is_nuclear)
        
        if content.get("notes"):
            try: slide.notes_slide.notes_text_frame.text = content["notes"]
            except: pass

    def _build_sterile_stack(self, slide, content, is_nuclear=False):
        sw, sh = self.prs.slide_width, self.prs.slide_height
        
        # THE BRANDING FORTRESS (2.2 Inch Top Offset)
        # Use detected floor if higher than 2.2", else default to 2.2" lock-down
        detected_floor = self.template_profile.get("brand_safe_floor")
        safe_top = Emu(Inches(2.2))
        
        # SPACE OPTIMIZATION: If slide is short (16:9) and has visuals, slightly compress guard to 1.8"
        has_visuals = bool(content.get("chart_data") or content.get("table_data") or content.get("image_prompt"))
        if sh < Emu(Inches(7.0)) and has_visuals:
            safe_top = Emu(Inches(1.8))
            
        if detected_floor and detected_floor > safe_top:
            safe_top = detected_floor
            
        safe_bottom = sh - Emu(Pt(20))
        
        title = content.get("title", "Untitled")
        subtitle = content.get("subtitle", "")
        bullets = [b.strip() for b in content.get("bullet_points", []) if b.strip()]
        chart_data = content.get("chart_data")
        table_data = content.get("table_data")

        # Auto-fit recursive loop
        font_scale = 1.0
        for attempt in range(5):
            y_curr = safe_top
            added_shapes = []

            if has_visuals:
                # VISUAL INSIGHT LAYOUT: Safe 3-line title + max canvas for chart/table
                t_size = int(self._theme["title_size"].pt * font_scale)
                # Allocate height for up to 3 wrapped lines to prevent title/visual overlap
                t_h = Emu(Pt(t_size * 3.5))
                t_box = slide.shapes.add_textbox(Emu(Inches(0.8)), y_curr, sw - Emu(Inches(1.6)), t_h)
                self._set_text(t_box, title or " ", Pt(t_size), True, self._theme["title_font"], self._theme["title_color"])
                y_curr += t_h
                added_shapes.append(t_box)

                # Give the full remaining height to the visual
                g_h = safe_bottom - y_curr
                if g_h > Emu(Pt(20)):
                    g_geom = (Emu(Inches(0.8)), y_curr, sw - Emu(Inches(1.6)), g_h)
                    v_shape = None
                    if chart_data:
                        v_shape = self._add_chart(slide, chart_data, g_geom)
                    elif table_data:
                        v_shape = self._add_table(slide, table_data, g_geom)
                    elif content.get("image_prompt"):
                        v_shape = self._add_image_marker(slide, content.get("image_prompt"), g_geom)

                    if v_shape:
                        added_shapes.append(v_shape)
                        y_curr += g_h

            else:
                # ============================================
                # NARRATIVE LAYOUT: The Goldilocks Zone (Balanced Pro)
                # ============================================
                title_gap = Emu(Pt(35))    
                subtitle_gap = Emu(Pt(25)) 

                # --- TITLE ---
                t_size = int(self._theme["title_size"].pt * font_scale)
                t_h = Emu(Pt(t_size * 1.6))
                t_box = slide.shapes.add_textbox(Emu(Inches(0.8)), y_curr, sw - Emu(Inches(1.6)), t_h)
                self._set_text(t_box, title or " ", Pt(t_size), True, self._theme["title_font"], self._theme["title_color"])
                
                if not subtitle:
                    y_curr += t_h + Emu(Pt(45))
                else:
                    y_curr += t_h + title_gap
                added_shapes.append(t_box)

                # --- SUBTITLE ---
                if subtitle:
                    st_size = int(self._theme["subtitle_size"].pt * font_scale)
                    st_h = Emu(Pt(st_size * 1.6))
                    st_box = slide.shapes.add_textbox(Emu(Inches(0.8)), y_curr, sw - Emu(Inches(1.6)), st_h)
                    self._set_text(st_box, subtitle or " ", Pt(st_size), False, self._theme["body_font"], self._theme["body_color"])
                    y_curr += st_h + subtitle_gap
                    added_shapes.append(st_box)

                # --- BULLETS ---
                b_size = int(self._theme["body_size"].pt * font_scale)
                display_bullets = bullets if bullets else []
                if display_bullets:
                    b_h = Emu(Pt(len(display_bullets) * b_size * 2.1))
                    b_box = slide.shapes.add_textbox(Emu(Inches(0.8)), y_curr, sw - Emu(Inches(1.6)), b_h)
                    self._set_bullets(b_box, display_bullets, Pt(b_size))
                    y_curr += b_h
                    added_shapes.append(b_box)


            if y_curr <= safe_bottom: break
            else:
                for s in added_shapes: 
                    try: slide.shapes._spTree.remove(s._sp)
                    except: pass
                
                # ONLY Nuclear Shred if it's a visual/blank slide (Happy State)
                if is_nuclear:
                    try: 
                        while len(slide.shapes) > 0:
                            slide.shapes._spTree.remove(slide.shapes[-1]._sp)
                    except: pass
                
                font_scale -= 0.1

    def _set_text(self, box, text, size, bold, font_name, color):
        tf = box.text_frame
        tf.word_wrap = True
        if not tf.paragraphs: tf.add_paragraph()
        p = tf.paragraphs[0]
        p.text = str(text)
        p.font.size, p.font.bold = size, bold
        p.font.name, p.font.color.rgb = font_name, color

    def _set_bullets(self, box, bullets, size):
        from pptx.util import Pt
        from pptx.oxml.ns import qn
        from lxml import etree
        tf = box.text_frame
        tf.word_wrap = True
        for i, b_text in enumerate(bullets):
            p = tf.paragraphs[i] if i < len(tf.paragraphs) else tf.add_paragraph()
            p.text = str(b_text)
            p.level = 0
            p.font.size, p.font.name = size, self._theme["body_font"]
            p.font.color.rgb = self._theme["body_color"]
            # Add space-before each bullet for breathing room (like user's reference)
            pPr = p._pPr
            if pPr is None:
                pPr = p._p.get_or_add_pPr()
            pPr.set('spcBef', str(int(Pt(14).emu / 127)))  # Generous separation

    def _add_chart(self, slide, chart_data, geom):
        """Adds a chart using standardized AI schema. Returns the shape."""
        try:
            from pptx.chart.data import CategoryChartData
            data = CategoryChartData()
            data.categories = chart_data.get("categories", ["A", "B", "C"])
            
            values = chart_data.get("values")
            if values:
                data.add_series(chart_data.get("title", "Metric"), [float(v) for v in values])
            elif chart_data.get("series"):
                for series in chart_data["series"]:
                    data.add_series(series["name"], [float(v) for v in series["values"]])
            
            chart_type = XL_CHART_TYPE.COLUMN_CLUSTERED
            ct = str(chart_data.get("type", "")).lower()
            if "line" in ct: chart_type = XL_CHART_TYPE.LINE
            elif "bar" in ct: chart_type = XL_CHART_TYPE.BAR_CLUSTERED
            elif "pie" in ct: chart_type = XL_CHART_TYPE.PIE

            # For pie charts, constrain to a square-ish centered area for visual clarity
            if chart_type == XL_CHART_TYPE.PIE:
                avail_h = geom[3]
                side = min(geom[2], avail_h)  # square side = min(width, height)
                left_offset = geom[0] + (geom[2] - side) // 2
                geom = (left_offset, geom[1], side, side)

            chart_shape = slide.shapes.add_chart(chart_type, geom[0], geom[1], geom[2], geom[3], data)
            chart_shape.chart.has_title = False
            return chart_shape
        except Exception as e:
            print(f"Chart Render Failed: {e}")
            return None

    def _add_table(self, slide, table_data, geom):
        """Adds a table, capped to fit within the available height. Returns the shape."""
        try:
            if isinstance(table_data, dict):
                headers = table_data.get("headers", [])
                rows = table_data.get("rows", [])
                full_table = [headers] + rows if headers else rows
            else:
                full_table = table_data

            if not full_table: return None

            avail_h = geom[3]  # total available height in EMU
            # Calculate safe row height (min 0.45", max to fit all rows)
            row_h_emu = Emu(Inches(0.45))
            # Cap rows to what fits in the available height
            max_rows = max(1, int(avail_h / row_h_emu))
            full_table = full_table[:max_rows]

            num_rows = len(full_table)
            num_cols = len(full_table[0]) if num_rows > 0 else 0
            if num_cols == 0: return None

            # Use the exact available height so the table fills the canvas
            shape = slide.shapes.add_table(num_rows, num_cols, geom[0], geom[1], geom[2], avail_h)
            table = shape.table

            # Set explicit equal row heights to prevent overflow
            exact_row_h = avail_h // num_rows
            for row in table.rows:
                row.height = exact_row_h

            for r in range(num_rows):
                for c in range(num_cols):
                    cell = table.cell(r, c)
                    cell.text = str(full_table[r][c])
                    tf = cell.text_frame
                    tf.word_wrap = True
                    for para in tf.paragraphs:
                        para.font.size = Pt(11)
                    if r == 0:
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = self._theme["title_color"]
                        p = tf.paragraphs[0]
                        p.font.bold, p.font.color.rgb = True, RGBColor(255, 255, 255)
            return shape
        except Exception as e:
            print(f"Table Render Failed: {e}")
            return None

    def _add_image_marker(self, slide, prompt, geom):
        try:
             from pptx.enum.shapes import MSO_SHAPE
             shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, geom[0], geom[1], geom[2], geom[3])
             shape.fill.solid()
             shape.fill.fore_color.rgb = RGBColor(245, 245, 248)
             shape.line.color.rgb = RGBColor(210, 210, 220)
             p = shape.text_frame.add_paragraph()
             p.text = f"[AI VISUAL INSIGHT]\n{prompt[:60]}..."
             p.font.size, p.font.color.rgb = Pt(12), RGBColor(120, 120, 140)
             p.alignment = PP_ALIGN.CENTER
        except: pass
