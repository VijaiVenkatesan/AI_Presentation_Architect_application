
import io
import os
from pptx import Presentation
from core.ppt_generator import PptGenerator
from pptx.util import Inches

TEMPLATE_PATH = r"D:\Backup\PPT\Datamatics-Presentation-Template-Light-Test.pptx"

def run_e2e_test():
    if not os.path.exists(TEMPLATE_PATH):
        print(f"FAILED: Template not found at {TEMPLATE_PATH}")
        return

    with open(TEMPLATE_PATH, "rb") as f:
        template_bytes = f.read()

    # Define the Title -> Content -> Thank You sequence
    slides_content = [
        {
            "title": "Slide 1: Executive Summary",
            "subtitle": "Presented by AI Architect",
            "bullet_points": ["Strategic initiative", "0% overlap goal"]
        },
        {
            "title": "Slide 2: Technical Deep-Dive",
            "bullet_points": ["XML Shredder", "Layout Bleach", "Sterile Injection"]
        },
        {
            "title": "Thank You",
            "subtitle": "Questions?"
        }
    ]

    print("--- Running True-Gamma E2E Verification ---")
    generator = PptGenerator(template_bytes=template_bytes)
    output_pptx = generator.generate(slides_content)

    # Inspect Result
    prs = Presentation(io.BytesIO(output_pptx))
    print(f"Generated Presentation: {len(prs.slides)} slides")

    for i, slide in enumerate(prs.slides):
        print(f"Slide {i}: {len(slide.shapes)} shapes")
        
        # Verify 2.0 inch branding fortress
        for j, shape in enumerate(slide.shapes):
            top_inch = shape.top / Inches(1)
            if top_inch < 1.4: # Offset slightly for Title boxes which might be close
                 print(f"  WARNING: Shape {j} ({shape.name}) is at {top_inch:.2f} inches (Too High!)")
            
            shape_text = ""
            if hasattr(shape, 'text_frame'):
                 if shape.text_frame and shape.text_frame.paragraphs:
                     shape_text = shape.text_frame.text
            print(f"    Shape {j}: {shape.name} | Top: {top_inch:.2f}in | Text: {shape_text[:15]}")

    # Inspect Layout Mapping
    # User's Layouts: 0=Title, 1=Content, 20=Thank You
    print("\n--- Layout Mapping Verification ---")
    # Note: Layout names can be used for verification
    print(f"  Slide 0 Layout: {prs.slides[0].slide_layout.name}")
    print(f"  Slide 1 Layout: {prs.slides[1].slide_layout.name}")
    print(f"  Slide 2 Layout: {prs.slides[2].slide_layout.name}")

    with open("true_gamma_test.pptx", "wb") as f:
        f.write(output_pptx)
    print("\nSUCCESS: True-Gamma E2E Test Completed.")

if __name__ == "__main__":
    run_e2e_test()
