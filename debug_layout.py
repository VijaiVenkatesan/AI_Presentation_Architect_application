
from pptx import Presentation
from pptx.util import Inches, Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os

def create_debug_slide(output_path, template_path=None):
    if template_path and os.path.exists(template_path):
        prs = Presentation(template_path)
    else:
        prs = Presentation()
        prs.slide_width = Emu(12192000) # 13.33"
        prs.slide_height = Emu(6858000) # 7.5"

    slide_layout = prs.slide_layouts[4] if len(prs.slide_layouts) > 4 else prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    
    # 1. Draw a 2.2-inch Guard Line
    guard_y = Inches(2.2)
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, guard_y, prs.slide_width, Emu(Pt(2))
    )
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(255, 0, 0) # Red
    line.line.visible = False
    
    # 2. Add a TextBox at exactly 2.2-inch
    txBox = slide.shapes.add_textbox(Inches(1), guard_y, Inches(8), Inches(1))
    tf = txBox.text_frame
    tf.text = "THIS TEXT STARTS AT EXACTLY 2.2 INCHES"
    p = tf.paragraphs[0]
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(0, 0, 255) # Blue
    
    # 3. Add a label
    label = slide.shapes.add_textbox(Inches(0.2), guard_y - Inches(0.3), Inches(2), Inches(0.3))
    label.text_frame.text = "2.2 Inch Guard Line"
    label.text_frame.paragraphs[0].font.size = Pt(10)

    prs.save(output_path)
    print(f"Debug PPTX saved to {output_path}")
    print(f"Slide Height: {prs.slide_height.inches:.2f}\"")
    print(f"Guard Line Y: {guard_y.inches:.2f}\"")

if __name__ == "__main__":
    out = r"C:\Users\vijai.v\Downloads\debug_guard_v1.pptx"
    # Try to use the user's template if we can find it, otherwise default
    template = r"d:\Antigravity_Projects\AI_Presentation_Architecture_App\assets\sample_template.pptx"
    create_debug_slide(out, template if os.path.exists(template) else None)
