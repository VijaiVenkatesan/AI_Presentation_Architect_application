
from pptx import Presentation
import io
import os

# 1. Create a dummy presentation with a problematic layout
prs = Presentation()
# Use layout index 1 (Title and Content) or others
for i, layout in enumerate(prs.slide_layouts):
    slide = prs.slides.add_slide(layout)
    print(f"Layout {i} ({layout.name}): {len(slide.shapes)} shapes")
    for j, shape in enumerate(slide.shapes):
        print(f"  Shape {j}: {shape.name} | Type: {shape.shape_type} | XML: {shape.element.tag}")

# 2. Test the Nuclear Shredder
print("\n--- Nuclear Shredder Test ---")
slide = prs.slides[0]
print(f"Before: {len(slide.shapes)} shapes")

# Nuclear Shredder logic
spTree = slide.shapes._spTree
for element in list(spTree):
    # Standard p:sp, p:pic, p:graphicFrame
    tag = element.tag
    if "sp" in tag or "pic" in tag or "graphicFrame" in tag:
        spTree.remove(element)

print(f"After: {len(slide.shapes)} shapes")
