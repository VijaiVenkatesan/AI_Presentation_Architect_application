"""
Generate a professional sample template PPTX for testing.
Run this script once to create assets/sample_template.pptx
"""

import os
import sys
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def create_sample_template():
    prs = Presentation()
    prs.slide_width = Emu(12192000)   # 13.33"
    prs.slide_height = Emu(6858000)   #  7.5"

    W = prs.slide_width
    H = prs.slide_height

    # ── Slide 1: Title Slide ─────────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    # Set background
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0x0a, 0x0a, 0x1a)

    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:  # Title
            ph.text = "Enterprise Presentation"
            for run in ph.text_frame.paragraphs[0].runs:
                run.font.size = Pt(40)
                run.font.bold = True
                run.font.color.rgb = RGBColor(0x00, 0xd4, 0xff)
                run.font.name = "Calibri"
        elif ph.placeholder_format.idx == 1:  # Subtitle
            ph.text = "Powered by AI — Your Company Name"
            for run in ph.text_frame.paragraphs[0].runs:
                run.font.size = Pt(20)
                run.font.color.rgb = RGBColor(0xa0, 0xa0, 0xb8)
                run.font.name = "Calibri"

    # ── Slide 2: Content Slide ───────────────────────────────────────────
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    bg2 = slide2.background
    fill2 = bg2.fill
    fill2.solid()
    fill2.fore_color.rgb = RGBColor(0x0d, 0x0d, 0x24)

    for ph in slide2.placeholders:
        if ph.placeholder_format.idx == 0:
            ph.text = "Key Highlights"
            for run in ph.text_frame.paragraphs[0].runs:
                run.font.size = Pt(32)
                run.font.bold = True
                run.font.color.rgb = RGBColor(0x00, 0xd4, 0xff)
                run.font.name = "Calibri"
        elif ph.placeholder_format.idx == 1:
            tf = ph.text_frame
            tf.clear()
            points = [
                "Strategic growth through AI adoption",
                "Revenue increase by 40% year-over-year",
                "Market expansion across 12 new regions",
                "Customer satisfaction score: 94%",
            ]
            for i, pt in enumerate(points):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = pt
                p.level = 0
                for run in p.runs:
                    run.font.size = Pt(16)
                    run.font.color.rgb = RGBColor(0xe0, 0xe0, 0xe8)
                    run.font.name = "Calibri"
                p.space_after = Pt(8)

    # ── Slide 3: Blank for charts ────────────────────────────────────────
    layout_idx = 6 if len(prs.slide_layouts) > 6 else len(prs.slide_layouts) - 1
    slide3 = prs.slides.add_slide(prs.slide_layouts[layout_idx])
    bg3 = slide3.background
    fill3 = bg3.fill
    fill3.solid()
    fill3.fore_color.rgb = RGBColor(0x11, 0x11, 0x28)

    # Add a title text box
    txb = slide3.shapes.add_textbox(
        Emu(int(W * 0.05)), Emu(int(H * 0.05)),
        Emu(int(W * 0.9)), Emu(int(H * 0.1))
    )
    p = txb.text_frame.paragraphs[0]
    p.text = "Data & Analytics"
    run = p.runs[0]
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x00, 0xd4, 0xff)
    run.font.name = "Calibri"

    # Save
    out_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), "assets")
    os.makedirs(out_dir, exist_ok=True)
    out_path = os.path.join(out_dir, "sample_template.pptx")
    prs.save(out_path)
    print(f"Sample template saved to: {out_path}")
    return out_path


if __name__ == "__main__":
    create_sample_template()
