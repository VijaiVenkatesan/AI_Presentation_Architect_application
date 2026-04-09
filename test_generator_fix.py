
import os
import sys
import io
from pptx import Presentation

# Add path
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))

from core.ppt_generator import PptGenerator
from core.template_parser import PptxTemplateParser

def test_generation():
    template_path = r"d:\Antigravity_Projects\AI_Presentation_Architecture_App\assets\sample_template.pptx"
    out_path = r"C:\Users\vijai.v\Downloads\test_fix_v1.pptx"
    
    # 1. Parse Template
    with open(template_path, "rb") as f:
        template_bytes = f.read()
    
    parser = PptxTemplateParser(template_bytes)
    profile = parser.parse()
    
    # 2. Generate Slides
    slides = [
        {
            "title": "Fix Verification: Slide 1",
            "subtitle": "No Repaired Warning Allowed",
            "bullet_points": ["Safe clearing implemented", "Branding guard locked at 2.2\""],
            "notes": "Speaker notes check"
        },
        {
            "title": "Fix Verification: Visual Slide",
            "bullet_points": [],
            "chart_data": {"categories": ["Q1", "Q2"], "series": [{"name": "Sales", "values": [100, 200]}]}
        }
    ]
    
    generator = PptGenerator(template_bytes=template_bytes, template_profile=profile)
    pptx_bytes = generator.generate(slides)
    
    with open(out_path, "wb") as f:
        f.write(pptx_bytes)
    
    print(f"Test PPTX saved to {out_path}")
    
    # 3. Verify Layouts
    try:
        prs = Presentation(out_path)
        print(f"File is readable. Slide count: {len(prs.slides)}")
        
        for i, slide in enumerate(prs.slides):
            idx = prs.slide_layouts.index(slide.slide_layout)
            name = slide.slide_layout.name
            print(f"Slide {i+1} Layout: {idx} ({name})")
            
            if i == 0:
                if idx != 0: print("❌ Slide 1 is NOT Layout 0 (Title)")
                else: print("✅ Slide 1 is Title Layout")
            elif i == len(prs.slides) - 1:
                if "thank" not in name.lower() and idx != 20: 
                    print(f"❌ Final Slide {i+1} is NOT Thank You layout (Current: {name})")
                else: print(f"✅ Final Slide is Thank You Layout")
            else:
                if "blank" not in name.lower() and idx != 4:
                    print(f"❌ Content Slide {i+1} is NOT Blank layout (Current: {name})")
                else: print(f"✅ Content Slide is Blank Layout")
                
    except Exception as e:
        print(f"Verification Error: {e}")

if __name__ == "__main__":
    test_generation()
