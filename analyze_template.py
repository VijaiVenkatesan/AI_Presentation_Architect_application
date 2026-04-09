
from pptx import Presentation
import os

TEMPLATE_PATH = r"D:\Backup\PPT\Datamatics-Presentation-Template-Light-Test.pptx"

if not os.path.exists(TEMPLATE_PATH):
    print(f"FAILED: File not found at {TEMPLATE_PATH}")
    exit(1)

prs = Presentation(TEMPLATE_PATH)

print(f"--- Template Analysis: {os.path.basename(TEMPLATE_PATH)} ---")
print(f"Masters: {len(prs.slide_masters)}")

for i, master in enumerate(prs.slide_masters):
    print(f"Master {i}: {len(master.shapes)} shapes")
    for j, shape in enumerate(master.shapes):
        print(f"  Shape {j}: {shape.name} | Type: {shape.shape_type}")
    
    print(f"  Layouts in Master {i}: {len(master.slide_layouts)}")
    for k, layout in enumerate(master.slide_layouts):
        print(f"    Layout {k} ({layout.name}): {len(layout.shapes)} shapes")
        for m, shape in enumerate(layout.shapes):
            # Check for types like PICTURE (13) or GRAPHIC_FRAME (14)
            print(f"      Shape {m}: {shape.name} | Type ID: {shape.shape_type} | Placeholder: {shape.is_placeholder}")
            if shape.is_placeholder:
                print(f"        Placeholder Index: {shape.placeholder_format.idx}")

print("\n--- Shredding Proof of Concept ---")
# Try to shred Layout 1 of Master 0
layout = prs.slide_masters[0].slide_layouts[1]
print(f"Testing Shredder on '{layout.name}'...")
spTree = layout.shapes._spTree
print(f"  Before: {len(list(spTree))} XML nodes")
for element in list(spTree):
    tag = element.tag
    if any(x in tag for x in ["sp", "pic", "graphicFrame"]):
        spTree.remove(element)
print(f"  After: {len(list(spTree))} XML nodes")
